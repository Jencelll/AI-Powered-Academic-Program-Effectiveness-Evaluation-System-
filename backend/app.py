# app.py
from flask import Flask, request, jsonify, redirect, url_for, g, send_file
from flask_cors import CORS # Enable CORS for frontend interaction
from models import db, FacultyUpload, SubjectAnalysis, AnalyticsData, IdempotencyKey, User, UniversalUploadLog, StudentRiskAssessment, StudentRiskUpload, StudentSubjectRecord, StudentAcademicProfile, SubjectAccomplishment, AuditLog
import os
import json # Import json to handle stored data
import requests # Used for utility functions
from datetime import datetime, timezone, timedelta
from functools import wraps
from werkzeug.security import generate_password_hash, check_password_hash
from werkzeug.utils import secure_filename
from itsdangerous import URLSafeTimedSerializer, BadSignature, SignatureExpired
from docx import Document
from docx.shared import Pt, Inches, RGBColor
from docx.enum.text import WD_ALIGN_PARAGRAPH
from io import BytesIO

app = Flask(__name__)
app.config['SECRET_KEY'] = 'your_secret_key_here' # Change this!
# Use SQLite in the project-level instance directory to persist historical data
project_root = os.path.abspath(os.path.join(os.path.dirname(__file__), '..'))
db_dir = os.path.join(project_root, 'instance')
os.makedirs(db_dir, exist_ok=True)
db_path = os.path.join(db_dir, 'cqi_data.db')
app.config['SQLALCHEMY_DATABASE_URI'] = 'sqlite:///' + db_path
app.config['SQLALCHEMY_TRACK_MODIFICATIONS'] = False

CORS(app) # Enable CORS for all routes
db.init_app(app)

# --- Auth helpers ---
serializer = URLSafeTimedSerializer(app.config['SECRET_KEY'])

def issue_token(user: User):
    payload = {'uid': user.id, 'role': user.role}
    return serializer.dumps(payload)

def parse_token(token: str):
    try:
        data = serializer.loads(token, max_age=60 * 60 * 12)  # 12 hours
        return data
    except SignatureExpired:
        return None
    except BadSignature:
        return None

def get_auth_token():
    auth = request.headers.get('Authorization', '')
    parts = auth.split()
    if len(parts) == 2 and parts[0].lower() == 'bearer':
        return parts[1]
    return None

def auth_required(f):
    @wraps(f)
    def wrapper(*args, **kwargs):
        token = get_auth_token()
        if not token:
            return jsonify({'error': 'Authentication required'}), 401
        data = parse_token(token)
        if not data:
            return jsonify({'error': 'Invalid or expired token'}), 401
        user = User.query.get(data.get('uid'))
        if not user or not user.active:
            return jsonify({'error': 'User not found or inactive'}), 401
        g.current_user = user
        return f(*args, **kwargs)
    return wrapper

def roles_required(*roles):
    def decorator(f):
        @wraps(f)
        def wrapper(*args, **kwargs):
            # Must be authenticated first
            token = get_auth_token()
            if not token:
                return jsonify({'error': 'Authentication required'}), 401
            data = parse_token(token)
            if not data:
                return jsonify({'error': 'Invalid or expired token'}), 401
            user = User.query.get(data.get('uid'))
            if not user or not user.active:
                return jsonify({'error': 'User not found or inactive'}), 401
            if user.role not in roles:
                return jsonify({'error': 'Access denied', 'required_roles': roles}), 403
            g.current_user = user
            return f(*args, **kwargs)
        return wrapper
    return decorator

@app.route('/api/admin/update-email', methods=['PUT'])
@roles_required('Admin')
def update_admin_email():
    data = request.get_json()
    current_password = data.get('current_password')
    new_email = data.get('new_email')

    if not current_password or not new_email:
        return jsonify({'error': 'Missing required fields'}), 400

    # Verify current password
    user = g.current_user
    if not check_password_hash(user.password_hash, current_password):
        return jsonify({'error': 'Incorrect current password'}), 401

    # Check if email is already taken
    if User.query.filter_by(email=new_email).first():
        return jsonify({'error': 'Email is already in use'}), 409

    # Update email
    old_email = user.email
    user.email = new_email
    db.session.commit()

    # Log action
    log_entry = AuditLog(
        user_id=user.id,
        action='email_change',
        details=json.dumps({'old_email': old_email, 'new_email': new_email}),
        ip_address=request.remote_addr,
        timestamp=datetime.utcnow()
    )
    db.session.add(log_entry)
    db.session.commit()

    return jsonify({'message': 'Email updated successfully'}), 200


@app.route('/api/admin/update-password', methods=['PUT'])
@roles_required('Admin')
def update_admin_password():
    data = request.get_json()
    current_password = data.get('current_password')
    new_password = data.get('new_password')

    if not current_password or not new_password:
        return jsonify({'error': 'Missing required fields'}), 400

    # Verify current password
    user = g.current_user
    if not check_password_hash(user.password_hash, current_password):
        return jsonify({'error': 'Incorrect current password'}), 401
    
    # Password policy check
    if len(new_password) < 8:
         return jsonify({'error': 'Password must be at least 8 characters long'}), 400
    
    import re
    if not re.search(r'[A-Z]', new_password):
        return jsonify({'error': 'Password must contain at least one uppercase letter'}), 400
    if not re.search(r'[a-z]', new_password):
        return jsonify({'error': 'Password must contain at least one lowercase letter'}), 400
    if not re.search(r'[0-9]', new_password):
        return jsonify({'error': 'Password must contain at least one number'}), 400
    if not re.search(r'[\W_]', new_password):
        return jsonify({'error': 'Password must contain at least one special character'}), 400

    if check_password_hash(user.password_hash, new_password):
        return jsonify({'error': 'New password cannot be the same as the current password'}), 400

    # Update password
    user.password_hash = generate_password_hash(new_password)
    db.session.commit()

    # Log action
    log_entry = AuditLog(
        user_id=user.id,
        action='password_change',
        details=json.dumps({'ip': request.remote_addr}),
        ip_address=request.remote_addr,
        timestamp=datetime.utcnow()
    )
    db.session.add(log_entry)
    db.session.commit()

    return jsonify({'message': 'Password updated successfully'}), 200


@app.route('/api/health', methods=['GET'])
def health_check():
    return jsonify({'status': 'ok', 'timestamp': to_manila_iso(datetime.utcnow())})



MANILA_TZ = timezone(timedelta(hours=8))
UTC_TZ = timezone.utc

def to_manila_iso(dt: datetime):
    if dt is None:
        return datetime.now(MANILA_TZ).isoformat()
    base = dt
    if base.tzinfo is None:
        base = base.replace(tzinfo=UTC_TZ)
    return base.astimezone(MANILA_TZ).isoformat()

# Ensure the database exists
with app.app_context():
    db.create_all()
    # Lightweight migration: add missing columns to 'user' table if not present
    try:
        from sqlalchemy import text
        conn = db.engine.connect()
        cols = conn.execute(text("PRAGMA table_info('user')")).fetchall()
        existing_cols = {row[1] for row in cols}
        alter_statements = []
        if 'full_name' not in existing_cols:
            alter_statements.append("ALTER TABLE user ADD COLUMN full_name VARCHAR(255)")
        if 'faculty_id' not in existing_cols:
            alter_statements.append("ALTER TABLE user ADD COLUMN faculty_id VARCHAR(64)")
        if 'student_id' not in existing_cols:
            alter_statements.append("ALTER TABLE user ADD COLUMN student_id VARCHAR(64)")
        if 'program' not in existing_cols:
            alter_statements.append("ALTER TABLE user ADD COLUMN program VARCHAR(120)")
        for stmt in alter_statements:
            try:
                conn.execute(text(stmt))
            except Exception:
                pass
        try:
            conn.execute(text("CREATE UNIQUE INDEX IF NOT EXISTS ix_user_faculty_id ON user(faculty_id)"))
        except Exception:
            pass
        try:
            conn.execute(text("CREATE UNIQUE INDEX IF NOT EXISTS ix_user_student_id ON user(student_id)"))
        except Exception:
            pass
        try:
            sru_cols = conn.execute(text("PRAGMA table_info('student_risk_upload')")).fetchall()
            sru_existing = {row[1] for row in sru_cols}
            if 'course' not in sru_existing:
                try:
                    conn.execute(text("ALTER TABLE student_risk_upload ADD COLUMN course VARCHAR(200)"))
                except Exception:
                    pass
            if 'year_level' not in sru_existing:
                try:
                    conn.execute(text("ALTER TABLE student_risk_upload ADD COLUMN year_level VARCHAR(32)"))
                except Exception:
                    pass
            if 'semester' not in sru_existing:
                try:
                    conn.execute(text("ALTER TABLE student_risk_upload ADD COLUMN semester VARCHAR(64)"))
                except Exception:
                    pass
            if 'section' not in sru_existing:
                try:
                    conn.execute(text("ALTER TABLE student_risk_upload ADD COLUMN section VARCHAR(64)"))
                except Exception:
                    pass
            if 'faculty_name' not in sru_existing:
                try:
                    conn.execute(text("ALTER TABLE student_risk_upload ADD COLUMN faculty_name VARCHAR(255)"))
                except Exception:
                    pass
            if 'uploaded_at' not in sru_existing:
                try:
                    conn.execute(text("ALTER TABLE student_risk_upload ADD COLUMN uploaded_at DATETIME"))
                except Exception:
                    pass
        except Exception:
            pass
        try:
            ssr_cols = conn.execute(text("PRAGMA table_info('student_subject_record')")).fetchall()
            ssr_existing = {row[1] for row in ssr_cols}
            if 'year_level' not in ssr_existing:
                try:
                    conn.execute(text("ALTER TABLE student_subject_record ADD COLUMN year_level VARCHAR(32)"))
                except Exception:
                    pass
        except Exception:
            pass
        try:
            sap_cols = conn.execute(text("PRAGMA table_info('student_academic_profile')")).fetchall()
            sap_existing = {row[1] for row in sap_cols}
            if 'failed_count' not in sap_existing:
                try:
                    conn.execute(text("ALTER TABLE student_academic_profile ADD COLUMN failed_count INTEGER DEFAULT 0"))
                except Exception:
                    pass
            if 'passed_count' not in sap_existing:
                try:
                    conn.execute(text("ALTER TABLE student_academic_profile ADD COLUMN passed_count INTEGER DEFAULT 0"))
                except Exception:
                    pass
            if 'risk_level' not in sap_existing:
                try:
                    conn.execute(text("ALTER TABLE student_academic_profile ADD COLUMN risk_level VARCHAR(32)"))
                except Exception:
                    pass
        except Exception:
            pass
        
        # Ensure academic_year exists on faculty_upload
        try:
            fu_cols = conn.execute(text("PRAGMA table_info('faculty_upload')")).fetchall()
            fu_existing = {row[1] for row in fu_cols}
            if 'academic_year' not in fu_existing:
                try:
                    conn.execute(text("ALTER TABLE faculty_upload ADD COLUMN academic_year VARCHAR(20)"))
                except Exception:
                    pass
        except Exception:
            pass

        # Ensure academic_year exists on student_risk_upload
        try:
            sru_cols = conn.execute(text("PRAGMA table_info('student_risk_upload')")).fetchall()
            sru_existing = {row[1] for row in sru_cols}
            if 'academic_year' not in sru_existing:
                try:
                    conn.execute(text("ALTER TABLE student_risk_upload ADD COLUMN academic_year VARCHAR(20)"))
                except Exception:
                    pass
        except Exception:
            pass

        # Ensure academic_year and semester exist on subject_accomplishment
        try:
            sa_cols = conn.execute(text("PRAGMA table_info('subject_accomplishment')")).fetchall()
            sa_existing = {row[1] for row in sa_cols}
            if 'academic_year' not in sa_existing:
                try:
                    conn.execute(text("ALTER TABLE subject_accomplishment ADD COLUMN academic_year VARCHAR(20)"))
                except Exception:
                    pass
            if 'semester' not in sa_existing:
                try:
                    conn.execute(text("ALTER TABLE subject_accomplishment ADD COLUMN semester VARCHAR(20)"))
                except Exception:
                    pass
        except Exception:
            pass

        conn.close()
    except Exception:
        # Non-fatal; continue if PRAGMA not available or engine not ready
        pass



# Import and run analysis function
from ai_processor import run_analysis

# --- Student Risk helpers ---
def _safe_float(x):
    try:
        import math
        if x is None:
            return None
        if isinstance(x, (int, float)):
            v = float(x)
            return None if math.isnan(v) else v
        s = str(x).strip()
        if not s:
            return None
        s = s.replace('%', '')
        v = float(s)
        return None if math.isnan(v) else v
    except Exception:
        return None

def _assign_risk(avg, failed, incomplete):
    avg = float(avg or 0)
    f = int(failed or 0)
    inc = int(incomplete or 0)
    if avg < 50 or f >= 4 or inc >= 4:
        return 'Critical'
    if avg < 60 or f >= 2 or inc >= 2:
        return 'High Risk'
    if avg > 80:
        return 'Low Risk'
    return 'Medium Risk'

def _recommendation(name, weakest_subjects):
    if not weakest_subjects:
        return 'Performance stable. Continue consistent study habits.'
    if len(weakest_subjects) == 1:
        return f'Needs intervention for {weakest_subjects[0]}.'
    return f'Needs intervention for {", ".join(weakest_subjects[:2])}.'

def _load_career_dataset():
    import pandas as pd
    global _career_cache, _career_cache_ts
    try:
        import time
        if _career_cache and _career_cache_ts and (time.time() - _career_cache_ts) < 600:
            return _career_cache
    except Exception:
        pass
    p1 = os.path.join(app.root_path, 'LSPU_Assessment_Data.xls')
    p2 = os.path.join(app.root_path, 'LSPU_Assessment_Data.xlsx')
    ds_path = p1 if os.path.exists(p1) else p2
    mappings = {
        'skills_by_course': {},
        'careers_by_skill': {},
        'job_paths_by_program': {},
        'employers_by_program': {},
    }
    try:
        xls = pd.ExcelFile(ds_path)
        for sheet in xls.sheet_names:
            df = xls.parse(sheet)
            cols = [str(c).strip().lower() for c in df.columns]
            # Course → Skill(s)
            if any('course' in c or 'subject' in c for c in cols) and any('skill' in c for c in cols):
                course_col = next((c for c in df.columns if 'course' in str(c).lower() or 'subject' in str(c).lower()), None)
                skill_col = next((c for c in df.columns if 'skill' in str(c).lower()), None)
                if course_col is not None and skill_col is not None:
                    for _, row in df.iterrows():
                        course = str(row.get(course_col) or '').strip().upper()
                        skills_raw = row.get(skill_col)
                        if not course:
                            continue
                        skills = []
                        try:
                            skills = [s.strip() for s in str(skills_raw or '').split(',') if s.strip()]
                            skills = [s for s in skills if any(ch.isalpha() for ch in s)]
                            skills = [s.title() for s in skills]
                        except Exception:
                            skills = []
                        if skills:
                            mappings['skills_by_course'][course] = list(set(mappings['skills_by_course'].get(course, []) + skills))
            # Skill → Career Path
            if any('skill' in c for c in cols) and any('career' in c or 'path' in c for c in cols):
                skill_col = next((c for c in df.columns if 'skill' in str(c).lower()), None)
                career_col = next((c for c in df.columns if 'career' in str(c).lower() or 'path' in str(c).lower()), None)
                if skill_col is not None and career_col is not None:
                    for _, row in df.iterrows():
                        skill = str(row.get(skill_col) or '').strip()
                        careers_raw = row.get(career_col)
                        if not skill:
                            continue
                        careers = []
                        try:
                            careers = [s.strip() for s in str(careers_raw or '').split(',') if s.strip()]
                            careers = [s for s in careers if any(ch.isalpha() for ch in s)]
                            careers = [s.title() for s in careers]
                        except Exception:
                            careers = []
                        if careers:
                            mappings['careers_by_skill'][skill] = list(set(mappings['careers_by_skill'].get(skill, []) + careers))
        try:
            df_alumni = xls.parse('Alumni_Survey')
            df_alumni.columns = [str(c).strip() for c in df_alumni.columns]
            df_tracer = None
            try:
                df_tracer = xls.parse('Graduate_Tracer')
            except Exception:
                try:
                    df_tracer = xls.parse('Graduate_Tracer_Data')
                except Exception:
                    df_tracer = None
            employed_vals = ['Employed full-time','Employed part-time','Self-employed','Working full-time','Working part-time','Self-employed','STARTUP']
            jobs = {}
            employers = {}
            for prog in ('BSIT','BSCS'):
                aj = []
                if 'program' in df_alumni.columns and 'employment_status' in df_alumni.columns and 'job_title' in df_alumni.columns:
                    aj = df_alumni[df_alumni['program'].astype(str).str.contains(prog, case=False, na=False) & df_alumni['employment_status'].astype(str).isin(employed_vals)]['job_title'].dropna().astype(str).tolist()
                tj = []
                if df_tracer is not None and 'program' in df_tracer.columns and 'employment_status' in df_tracer.columns and 'job_title' in df_tracer.columns:
                    tj = df_tracer[df_tracer['program'].astype(str).str.contains(prog, case=False, na=False) & df_tracer['employment_status'].astype(str).isin(employed_vals)]['job_title'].dropna().astype(str).tolist()
                allj = [s.strip().title() for s in (aj + tj) if s and s.lower() not in ('nan','n/a','na')]
                if allj:
                    from collections import Counter
                    cnt = Counter(allj)
                    top = [k for k, _ in sorted(cnt.items(), key=lambda x: (-x[1], x[0]))][:5]
                    jobs[prog] = top
                # Employers/companies if available
                employer_cols = [c for c in df_alumni.columns if any(k in str(c).lower() for k in ('employer','company','organization','workplace'))]
                employer_vals = []
                if employer_cols:
                    try:
                        df_emp = df_alumni[df_alumni['program'].astype(str).str.contains(prog, case=False, na=False)]
                        for col in employer_cols:
                            employer_vals += df_emp[col].dropna().astype(str).tolist()
                    except Exception:
                        pass
                if df_tracer is not None:
                    tracer_employer_cols = [c for c in df_tracer.columns if any(k in str(c).lower() for k in ('employer','company','organization','workplace'))]
                    if tracer_employer_cols:
                        try:
                            df_temp = df_tracer[df_tracer['program'].astype(str).str.contains(prog, case=False, na=False)]
                            for col in tracer_employer_cols:
                                employer_vals += df_temp[col].dropna().astype(str).tolist()
                        except Exception:
                            pass
                employer_vals = [s.strip().title() for s in employer_vals if s and s.lower() not in ('nan','n/a','na')]
                if employer_vals:
                    from collections import Counter
                    ecnt = Counter(employer_vals)
                    etop = [k for k, _ in sorted(ecnt.items(), key=lambda x: (-x[1], x[0]))][:5]
                    employers[prog] = etop
            mappings['job_paths_by_program'] = jobs
            mappings['employers_by_program'] = employers
        except Exception:
            pass
    except Exception:
        pass
    base_skills_by_course = {
        'ALGORITHM': ['Analytical', 'Programming'],
        'DATA STRUCTURE': ['Analytical', 'Programming'],
        'SOFTWARE ENGINEERING': ['Design', 'Programming', 'Collaboration'],
        'DATABASE': ['Data', 'Web'],
        'WEB': ['Web', 'Design'],
        'NETWORK': ['Hardware', 'Infrastructure'],
        'SECURITY': ['Security', 'Research'],
        'HARDWARE': ['Hardware'],
    }
    base_careers_by_skill = {
        'Programming': ['Software Engineer', 'Systems Developer', 'Backend Engineer'],
        'Web': ['Frontend Developer', 'Full-stack Developer'],
        'Data': ['Data Analyst', 'Database Administrator'],
        'Infrastructure': ['IT Infrastructure', 'Network Administrator'],
        'Security': ['Security Analyst', 'Cybersecurity Engineer'],
        'Design': ['UI/UX Designer', 'Product Engineer'],
        'Analytical': ['Algorithm Engineer', 'Research Assistant'],
        'Hardware': ['Embedded Systems Engineer', 'Network Engineer'],
    }
    for k, v in base_skills_by_course.items():
        mappings['skills_by_course'][k] = list(set(mappings['skills_by_course'].get(k, []) + v))
    for k, v in base_careers_by_skill.items():
        mappings['careers_by_skill'][k] = list(set(mappings['careers_by_skill'].get(k, []) + v))
    try:
        import time
        _career_cache = mappings
        _career_cache_ts = time.time()
    except Exception:
        pass
    return mappings

def _career_insight_for_student(breakdown, mappings):
    # breakdown: list of {subject, grade, status}
    # Compute strong and weak subjects with grades
    grades = []
    for item in breakdown:
        gnum = _safe_float(item.get('grade'))
        if gnum is None:
            continue
        grades.append({'subject': item.get('subject') or '', 'grade': float(gnum), 'status': str(item.get('status') or '').upper()})
    strong_pairs = [it for it in grades if (it['grade'] >= 90 and it['status'] not in ('FAILED','INC','INCOMPLETE','W','WITHDRAWN'))]
    strong = [it['subject'] for it in strong_pairs]
    weak_pairs = sorted(grades, key=lambda x: (x['grade'], x['subject']))[:3]
    avg_val = 0.0
    if grades:
        try:
            avg_val = round(sum(it['grade'] for it in grades)/len(grades), 2)
        except Exception:
            avg_val = 0.0
    skills_counter = {}
    for subj in strong:
        key = str(subj or '').upper()
        matched = []
        for course_key, skills in mappings.get('skills_by_course', {}).items():
            if course_key in key:
                matched += skills
        for s in matched:
            if any(ch.isalpha() for ch in str(s)):
                skills_counter[s] = skills_counter.get(s, 0) + 1
    top_skills = sorted(skills_counter.items(), key=lambda x: (-x[1], x[0]))[:5]
    if not top_skills:
        guess = []
        for subj in strong:
            k = str(subj or '').upper()
            if 'WEB' in k:
                guess.append('Web')
            if 'DATA' in k or 'DB' in k or 'DATABASE' in k:
                guess.append('Data')
            if 'NET' in k or 'NETWORK' in k:
                guess.append('Infrastructure')
            if 'SEC' in k or 'SECURITY' in k:
                guess.append('Security')
            if 'ALGO' in k or 'ALGORITHM' in k:
                guess.append('Analytical')
                guess.append('Programming')
            if 'CMSC' in k or 'CSST' in k or 'IT' in k:
                guess.append('Programming')
        if not guess:
            guess = ['Programming', 'Analytical']
        from collections import Counter
        cnt = Counter(guess)
        top_skills = sorted(cnt.items(), key=lambda x: (-x[1], x[0]))[:5]
    career_counter = {}
    for skill, _ in top_skills:
        for c in mappings.get('careers_by_skill', {}).get(skill, []):
            career_counter[c] = career_counter.get(c, 0) + 1
    prog_guess = None
    keys = [str(s or '').upper() for s in strong]
    if any(('CMSC' in k) for k in keys):
        prog_guess = 'BSCS'
    if any((('ITST' in k) or ('ITEC' in k) or ('ITE' in k) or ('CSST' in k))) and prog_guess is None:
        prog_guess = 'BSIT'
    jp = mappings.get('job_paths_by_program', {}).get(prog_guess or '', [])
    for c in jp:
        career_counter[c] = career_counter.get(c, 0) + 2
    if not career_counter:
        for c in mappings.get('careers_by_skill', {}).get('Programming', []):
            career_counter[c] = career_counter.get(c, 0) + 1
    top_careers = sorted(career_counter.items(), key=lambda x: (-x[1], x[0]))[:5]
    if not top_careers:
        alts = mappings.get('job_paths_by_program', {}).get(prog_guess or 'BSIT', [])
        if not alts:
            alts = mappings.get('careers_by_skill', {}).get('Programming', [])
        top_careers = [(c, 1) for c in alts[:5]]

    skill_weights = {}
    for item in breakdown:
        g = _safe_float(item.get('grade')) or 0
        status = str(item.get('status') or '').upper()
        subj = str(item.get('subject') or '').upper()
        w = 0
        if status in ('FAILED', 'INC', 'INCOMPLETE', 'W', 'WITHDRAWN'):
            w = -1
        elif g >= 90:
            w = 3
        elif g >= 85:
            w = 2
        elif g >= 75:
            w = 1
        if w == 0:
            continue
        for course_key, skills in mappings.get('skills_by_course', {}).items():
            if course_key in subj:
                for s in skills:
                    if any(ch.isalpha() for ch in str(s)):
                        skill_weights[s] = skill_weights.get(s, 0) + w
    potential_counter = {}
    for skill, w in skill_weights.items():
        for c in mappings.get('careers_by_skill', {}).get(skill, []):
            potential_counter[c] = potential_counter.get(c, 0) + w
    for c in jp:
        potential_counter[c] = potential_counter.get(c, 0) + 2
    potential_sorted = sorted(potential_counter.items(), key=lambda x: (-x[1], x[0]))[:8]
    if not potential_sorted:
        fallback = mappings.get('careers_by_skill', {}).get('Programming', [])
        potential_sorted = [(c, 1) for c in fallback[:8]]
    # Infer program from all subjects if strong was empty
    if prog_guess is None:
        keys_all = [str((it or {}).get('subject') or '').upper() for it in breakdown]
        if any(('CMSC' in k) for k in keys_all):
            prog_guess = 'BSCS'
        elif any((('ITST' in k) or ('ITEC' in k) or ('ITE' in k) or ('CSST' in k))) and prog_guess is None:
            prog_guess = 'BSIT'
    workplaces = mappings.get('employers_by_program', {}).get(prog_guess or '', [])
    strong_subjects = [{'subject': it['subject'], 'grade': it['grade']} for it in sorted(strong_pairs, key=lambda x: (-x['grade'], x['subject']))[:3]]
    weak_subjects = [{'subject': it['subject'], 'grade': it['grade']} for it in weak_pairs]
    top_weighted = sorted(skill_weights.items(), key=lambda x: (-x[1], x[0]))[:3]
    skills_str = ", ".join([f"{s}(+{w})" if w > 0 else f"{s}({w})" for s, w in top_weighted])
    failures = [str((it or {}).get('subject') or '') for it in breakdown if str((it or {}).get('status') or '').upper() in ('FAILED','INC','INCOMPLETE','W','WITHDRAWN')]
    rsn_parts = []
    if strong_subjects:
        rsn_parts.append("High scores in " + ", ".join([f"{it['subject']}({it['grade']})" for it in strong_subjects]))
    else:
        rsn_parts.append("Inferred strengths from overall subject grades")
    if skills_str:
        rsn_parts.append(f"Top skills: {skills_str}")
    if prog_guess:
        rsn_parts.append(f"Program context: {prog_guess}")
    if workplaces:
        rsn_parts.append(f"Common workplaces: {', '.join(workplaces[:3])}")
    if failures:
        rsn_parts.append(f"Penalties due to {', '.join(failures[:2])}")
    reasoning = "; ".join([p for p in rsn_parts if p])
    # Build per-path reasons referencing subject and grade
    subj_to_domains = {}
    for item in breakdown:
        subj = str(item.get('subject') or '')
        key = subj.upper()
        domains = []
        for course_key, skills in mappings.get('skills_by_course', {}).items():
            if course_key in key:
                domains += skills
        subj_to_domains[subj] = domains
    strong_for_reason = strong_subjects if strong_subjects else [{'subject': it['subject'], 'grade': it['grade']} for it in sorted(grades, key=lambda x: (-x['grade'], x['subject']))[:3]]
    path_reasons = []
    for name, sc in top_careers[:3]:
        matched = None
        for it in strong_for_reason:
            ds = [d for d in subj_to_domains.get(it['subject'], []) if any(ch.isalpha() for ch in str(d))]
            if ds:
                matched = it
                break
        if matched is None and strong_for_reason:
            matched = strong_for_reason[0]
        reason = None
        if matched:
            dom = subj_to_domains.get(matched['subject'], [])
            dom_str = (dom[0] if dom else 'domain')
            reason = f"based on high grade in {matched['subject']} ({matched['grade']}) and strength in {dom_str}"
        else:
            reason = "based on relative performance pattern across subjects"
        path_reasons.append({'name': name, 'reason': reason, 'score': sc})
    # Support tracks if overall low
    support_tracks = []
    low_overall = (avg_val < 75.0) or (len([1 for it in grades if it['status'] in ('FAILED','INC','INCOMPLETE','W','WITHDRAWN')]) >= 2)
    if low_overall:
        support_tracks = [
            'Bridge Training: Programming Fundamentals',
            'Remedial: Networking Basics',
            'Study Plan: Data Handling & Databases'
        ]
    return {
        'strong_skill_clusters': [s for s, _ in top_skills],
        'recommended_paths': [{'name': c, 'score': sc} for c, sc in top_careers],
        'potential_paths': [{'name': c, 'score': sc} for c, sc in potential_sorted],
        'workplaces': workplaces,
        'subjects': strong,
        'strong_subjects': strong_subjects,
        'weak_subjects': weak_subjects,
        'reasoning': reasoning,
        'advice': 'Strengthen weaker areas with guided study plans.',
        'path_reasons': path_reasons,
        'average': avg_val,
        'support_tracks': support_tracks,
    }

def _extract_course_code(filepath: str):
    try:
        name = os.path.basename(filepath)
        base = os.path.splitext(name)[0]
        s = base.replace('_', ' ').upper()
        import re
        m = re.search(r"([A-Z]{2,}\s*\d{3,})", s)
        if m:
            return m.group(1).strip()
        return base.upper()
    except Exception:
        return None

def _extract_year_level(filepath: str, df=None):
    try:
        name = os.path.basename(filepath)
        base = os.path.splitext(name)[0]
        s = base.replace('_', ' ').upper()
        import re
        m = re.search(r"\b([1-4])\s*[A-Z]?\b", s)
        if m:
            return m.group(1)
        m2 = re.search(r"BS\w+\s+([IVX]{1,3})[AB]?", s)
        if m2:
            roman = m2.group(1)
            mapping = {'I': '1', 'II': '2', 'III': '3', 'IV': '4'}
            return mapping.get(roman, None)
        if df is not None:
            try:
                cols = [str(c).strip().upper() for c in df.columns]
                for k in ('PROGRAM','SECTION'):
                    if k in cols:
                        col = next(c for c in df.columns if str(c).strip().upper()==k)
                        vals = df[col].astype(str).str.upper().head(10).tolist()
                        for v in vals:
                            m = re.search(r"\b([1-4])\s*[A-Z]?\b", v)
                            if m:
                                return m.group(1)
                            m2 = re.search(r"BS\w+\s+([IVX]{1,3})[AB]?", v)
                            if m2:
                                roman = m2.group(1)
                                mapping = {'I': '1', 'II': '2', 'III': '3', 'IV': '4'}
                                return mapping.get(roman, None)
            except Exception:
                pass
        return None
    except Exception:
        return None

def _is_valid_name(value: str):
    s = str(value or '').strip()
    if not s:
        return False
    up = s.upper()
    bad = (
        'STUDENT', 'TOTAL', 'AVERAGE', 'SUMMARY', 'SCHEDULE', 'PREPARED',
        'CERTIFIED', 'CHECKED', 'SUBMITTED', 'INSTRUCTOR', 'DEPARTMENT',
        'NOTED', 'NOTE', 'BOHOL'
    )
    if any(k in up for k in bad):
        return False
    if ':' in s:
        return False
    try:
        import re
        if re.search(r"\d", s):
            return False
        words = re.split(r"[,\s]+", s)
        words = [w for w in words if w]
        if len(words) < 2:
            return False
        letters = sum(ch.isalpha() for ch in s)
        if letters < 4:
            return False
        return True
    except Exception:
        return False

def _clean_name(value: str):
    s = str(value or '').strip()
    try:
        import re
        s = s.replace('\u00a0', ' ')
        s = re.sub(r"\s+", " ", s)
        s = s.strip(' ,')
        return s.title()
    except Exception:
        return s

def _read_excel_sheet(filepath: str, header=None):
    import pandas as pd
    try:
        df = pd.read_excel(filepath, sheet_name='Summary', header=header)
        setattr(df, '_sheet_name', 'Summary')
        return df
    except Exception:
        xls = pd.ExcelFile(filepath)
        target = None
        for nm in xls.sheet_names:
            if 'summary' in str(nm).lower():
                target = nm
                break
        if target is None:
            target = xls.sheet_names[0]
        df = pd.read_excel(filepath, sheet_name=target, header=header)
        setattr(df, '_sheet_name', target)
        return df

def _parse_summary_excel(filepath):
    import pandas as pd
    raw = _read_excel_sheet(filepath, header=None)
    header_idx = None
    for i in range(min(len(raw), 30)):
        row_vals = [str(v).strip().upper() for v in list(raw.iloc[i].values)]
        if any('STUDENT' in v and 'NAME' in v for v in row_vals):
            header_idx = i
            break
    if header_idx is None:
        df = _read_excel_sheet(filepath, header=0)
        sheet_name = getattr(df, '_sheet_name', None) or 'Summary'
        header_row = 0
        rows = []
        warnings = []
        rows_skipped = 0
        cols = list(df.columns)
        name_col = next((c for c in df.columns if 'name' in str(c).lower() or 'student' in str(c).lower()), df.columns[0])
        def _looks_like_numeric_series(series):
            try:
                vals = [str(v).strip() for v in list(series.head(10).values)]
                hits = 0
                for v in vals:
                    if not v:
                        continue
                    try:
                        float(str(v).replace('%',''))
                        hits += 1
                    except Exception:
                        pass
                return hits >= max(3, len(vals)//2)
            except Exception:
                return False
        def _is_subject_col(col):
            key = str(col).strip().upper().replace('\n',' ')
            import re
            code_like = re.search(r"[A-Z]{2,}\s*\d{2,3}", key) is not None
            numeric_like = _looks_like_numeric_series(df[col])
            non_subject_keys = ('REMARK','REMARKS','STATUS','STUDENT','NAME','GENDER','ID')
            if any(k in key for k in non_subject_keys):
                return False
            return code_like or numeric_like
        subject_cols = [c for c in df.columns if _is_subject_col(c)]
        if not subject_cols:
            grade_col = None
            for key in ('equivalent','final grade','rounded off','final','finals','grade','average','midterm'):
                m = next((c for c in df.columns if key == str(c).strip().lower()), None)
                if m is not None:
                    grade_col = m
                    break
            if grade_col is None:
                grade_col = df.columns[-1]
            subject_cols = [grade_col]
        for ridx, row in df.iterrows():
            name_val = row.get(name_col)
            try:
                is_na = pd.isna(name_val)
            except Exception:
                is_na = False
            if is_na:
                rows_skipped += 1
                continue
            name_raw = str(name_val or '').strip()
            if not _is_valid_name(name_raw):
                rows_skipped += 1
                continue
            name = _clean_name(name_raw)
            breakdown = []
            for col in subject_cols:
                raw_val = row.get(col)
                norm = _safe_float(raw_val)
                invalid = (norm is None) or (float(norm or 0) < 0 or float(norm or 0) > 100)
                sval = ''
                try:
                    sval = str(raw_val).strip().upper()
                except Exception:
                    sval = ''
                if any(k in sval for k in ('INC','INCOMPLETE','W','WITHDRAWN')):
                    status = 'INC'
                else:
                    status = 'PASSED' if float(norm or 0) >= 75 else ('INVALID' if invalid else 'FAILED') if (norm is not None) else 'INVALID'
                if invalid:
                    warnings.append({'type': 'INVALID_GRADE', 'student': name, 'subject': str(col), 'raw': raw_val})
                breakdown.append({
                    'subject': str(col),
                    'header': str(col),
                    'raw': raw_val,
                    'grade': (norm if norm is not None else 0.0),
                    'normalized': norm,
                    'invalid': invalid,
                    'status': status,
                })
            rows.append({'name': name, 'row_index': int(ridx)+header_row+1, 'breakdown': breakdown})
        report = {
            'file_name': os.path.basename(filepath),
            'sheet_name': sheet_name,
            'header_row_index': header_row,
            'rows_parsed': len(rows),
            'rows_skipped': rows_skipped,
            'subject_columns': [str(c) for c in subject_cols],
            'warnings': warnings,
        }
        return rows, report
    header = list(raw.iloc[header_idx].values)
    df = raw.iloc[header_idx+1:].copy()
    df.columns = header
    cols_norm = {str(c).strip().upper(): c for c in df.columns}
    name_cand = next((cols_norm[k] for k in cols_norm.keys() if 'STUDENT' in k and 'NAME' in k), df.columns[0])
    def _looks_like_numeric_series(series):
        try:
            vals = [str(v).strip() for v in list(series.head(10).values)]
            hits = 0
            for v in vals:
                if not v:
                    continue
                try:
                    float(str(v).replace('%',''))
                    hits += 1
                except Exception:
                    pass
            return hits >= max(3, len(vals)//2)
        except Exception:
            return False
    cols = list(df.columns)
    try:
        name_idx = cols.index(name_cand)
    except Exception:
        name_idx = 0
    if _looks_like_numeric_series(df[name_cand]):
        try:
            if name_idx+1 < len(cols):
                alt = cols[name_idx+1]
                s = df.iloc[:, name_idx+1].astype(str).str.strip().head(10).tolist()
                if any(sum(ch.isalpha() for ch in v) >= 3 and (' ' in v or ',' in v) for v in s):
                    name_cand = alt
                    name_idx = name_idx+1
        except Exception:
            pass
    grade_cand = None
    for orig in df.columns:
        key = str(orig).strip().upper().replace('\n',' ')
        if ('FINAL' in key and 'GRADE' in key) or ('ROUNDED' in key and 'OFF' in key):
            grade_cand = orig
            break
    if grade_cand is None:
        for key in ('FINALS','FINAL','GRADE','AVERAGE','MIDTERM'):
            if key in cols_norm:
                grade_cand = cols_norm[key]
                break
    if grade_cand is None and 'EQUIVALENT' in cols_norm:
        grade_cand = cols_norm['EQUIVALENT']
    status_cand = None
    for key in ('REMARKS','REMARK','STATUS'):
        if key in cols_norm:
            status_cand = cols_norm[key]
            break
    subj = _extract_course_code(filepath) or 'COURSE'
    rows = []
    warnings = []
    sheet_name = getattr(df, '_sheet_name', None) or 'Summary'
    header_row = header_idx
    rows_skipped = 0
    # Detect subject columns: headers that look like subject codes or numeric columns
    def _is_subject_col(col):
        key = str(col).strip().upper().replace('\n',' ')
        import re
        code_like = re.search(r"[A-Z]{2,}\s*\d{2,3}", key) is not None
        numeric_like = _looks_like_numeric_series(df[col])
        non_subject_keys = ('REMARK','REMARKS','STATUS','STUDENT','NAME','GENDER','ID')
        if any(k in key for k in non_subject_keys):
            return False
        return code_like or numeric_like
    subject_cols = [c for c in df.columns if _is_subject_col(c)]
    for ridx, row in df.iterrows():
        try:
            name_val = row.get(name_cand)
        except Exception:
            name_val = None
        try:
            if hasattr(name_val, 'tolist'):
                lst = [v for v in name_val.tolist() if not (isinstance(v, float) and pd.isna(v))]
                name_val = (lst[0] if lst else None)
        except Exception:
            pass
        if (name_val is None or (isinstance(name_val, float) and pd.isna(name_val))) and (name_idx is not None):
            try:
                name_val = row.iloc[name_idx]
            except Exception:
                name_val = None
        try:
            is_na = pd.isna(name_val)
        except Exception:
            is_na = False
        if is_na:
            rows_skipped += 1
            continue
        name_raw = str(name_val or '').strip()
        if not _is_valid_name(name_raw):
            rows_skipped += 1
            continue
        name = _clean_name(name_raw)
        breakdown = []
        for col in subject_cols:
            raw_val = row.get(col)
            norm = _safe_float(raw_val)
            invalid = (norm is None) or (float(norm or 0) < 0 or float(norm or 0) > 100)
            status_val = ''
            sval = None
            try:
                sval = str(raw_val).strip().upper()
            except Exception:
                sval = ''
            if any(k in sval for k in ('INC','INCOMPLETE','W','WITHDRAWN')):
                status = 'INC'
            else:
                status = 'PASSED' if float(norm or 0) >= 75 else ('INVALID' if invalid else 'FAILED') if (norm is not None) else 'INVALID'
            if invalid:
                warnings.append({'type': 'INVALID_GRADE', 'student': name, 'subject': str(col), 'raw': raw_val})
            breakdown.append({
                'subject': str(col),
                'header': str(col),
                'raw': raw_val,
                'grade': (norm if norm is not None else 0.0),
                'normalized': norm,
                'invalid': invalid,
                'status': status,
            })
        rows.append({'name': name, 'row_index': int(ridx)+header_idx+1, 'breakdown': breakdown})
    report = {
        'file_name': os.path.basename(filepath),
        'sheet_name': sheet_name,
        'header_row_index': header_row,
        'rows_parsed': len(rows),
        'rows_skipped': rows_skipped,
        'subject_columns': [str(c) for c in subject_cols],
        'warnings': warnings,
    }
    return rows, report

def _parse_summary_for_risk_tracking(filepath):
    import pandas as pd
    raw = _read_excel_sheet(filepath, header=None)
    header_idx = None
    for i in range(min(len(raw), 30)):
        row_vals = [str(v).strip().upper() for v in list(raw.iloc[i].values)]
        if any('STUDENT' in v and 'NAME' in v for v in row_vals):
            header_idx = i
            break
    subject = _extract_course_code(filepath) or 'COURSE'
    if header_idx is None:
        df = _read_excel_sheet(filepath, header=0)
        year_level = _extract_year_level(filepath, df) or None
        name_col = next((c for c in df.columns if 'name' in str(c).lower()), df.columns[0])
        grade_col = None
        for key in ('equivalent','final grade','rounded off','final','grade','average','midterm'):
            m = next((c for c in df.columns if key == str(c).strip().lower()), None)
            if m is not None:
                grade_col = m
                break
        status_col = next((c for c in df.columns if 'remark' in str(c).lower() or 'status' in str(c).lower()), None)
        rows = []
        for _, row in df.iterrows():
            name_val = row.get(name_col)
            if pd.isna(name_val):
                continue
            name_raw = str(name_val or '').strip()
            if not _is_valid_name(name_raw):
                continue
            name = _clean_name(name_raw)
            gsrc = row.get(grade_col) if grade_col is not None else row.get(df.columns[-1])
            gv = _safe_float(gsrc)
            if gv is None:
                gv = 0.0
            status_val = ''
            if status_col:
                sval = row.get(status_col)
                if not pd.isna(sval):
                    status_val = str(sval or '').strip().upper()
            if not status_val:
                status_val = 'PASSED' if float(gv) >= 75 else 'FAILED'
            rows.append({'student_name': name, 'grade': gv, 'status': status_val, 'subject': subject, 'year_level': year_level})
        return rows
    header = list(raw.iloc[header_idx].values)
    df = raw.iloc[header_idx+1:].copy()
    df.columns = header
    cols_norm = {str(c).strip().upper(): c for c in df.columns}
    name_cand = next((cols_norm[k] for k in cols_norm.keys() if 'STUDENT' in k and 'NAME' in k), df.columns[0])
    cols = list(df.columns)
    try:
        name_idx = cols.index(name_cand)
    except Exception:
        name_idx = 0
    def _looks_like_numeric_series(series):
        try:
            vals = [str(v).strip() for v in list(series.head(10).values)]
            hits = 0
            for v in vals:
                if not v:
                    continue
                try:
                    float(str(v).replace('%',''))
                    hits += 1
                except Exception:
                    pass
            return hits >= max(3, len(vals)//2)
        except Exception:
            return False
    if _looks_like_numeric_series(df[name_cand]):
        try:
            if name_idx+1 < len(cols):
                alt = cols[name_idx+1]
                s = df.iloc[:, name_idx+1].astype(str).str.strip().head(10).tolist()
                if any(sum(ch.isalpha() for ch in v) >= 3 and (' ' in v or ',' in v) for v in s):
                    name_cand = alt
                    name_idx = name_idx+1
        except Exception:
            pass
    grade_cand = None
    for orig in df.columns:
        key = str(orig).strip().upper().replace('\n',' ')
        if ('FINAL' in key and 'GRADE' in key) or ('ROUNDED' in key and 'OFF' in key):
            grade_cand = orig
            break
    if grade_cand is None:
        for key in ('FINALS','FINAL','GRADE','AVERAGE','MIDTERM'):
            if key in cols_norm:
                grade_cand = cols_norm[key]
                break
    if grade_cand is None and 'EQUIVALENT' in cols_norm:
        grade_cand = cols_norm['EQUIVALENT']
    status_cand = None
    for key in ('REMARKS','REMARK','STATUS'):
        if key in cols_norm:
            status_cand = cols_norm[key]
            break
    year_level = _extract_year_level(filepath, df) or None
    rows = []
    for _, row in df.iterrows():
        try:
            name_val = row.get(name_cand)
        except Exception:
            name_val = None
        try:
            if hasattr(name_val, 'tolist'):
                lst = [v for v in name_val.tolist() if not (isinstance(v, float) and pd.isna(v))]
                name_val = (lst[0] if lst else None)
        except Exception:
            pass
        is_na = False
        if name_val is None:
            is_na = True
        elif isinstance(name_val, float):
            try:
                is_na = pd.isna(name_val)
            except Exception:
                is_na = False
        elif isinstance(name_val, str):
            is_na = (len(name_val.strip()) == 0)
        if is_na:
            continue
        name_raw = str(name_val or '').strip()
        if not _is_valid_name(name_raw):
            continue
        name = _clean_name(name_raw)
        gsrc = row.get(grade_cand) if grade_cand is not None else None
        gv = _safe_float(gsrc)
        if gv is None:
            alt = row.get(cols_norm.get('EQUIVALENT'))
            gv = _safe_float(alt) or 0.0
        status_val = ''
        if status_cand:
            sval = row.get(status_cand)
            if not pd.isna(sval):
                status_val = str(sval or '').strip().upper()
        if not status_val:
            status_val = 'PASSED' if float(gv) >= 75 else 'FAILED'
        rows.append({'student_name': name, 'grade': gv, 'status': status_val, 'subject': subject, 'year_level': year_level})
    return rows

@app.route('/api/student-risk/evaluate', methods=['POST'])
@roles_required('Admin', 'Faculty')
def evaluate_student_risk():
    try:
        if 'grades' not in request.files:
            return jsonify({'error': 'Missing file field: grades'}), 400
        f = request.files['grades']
        upload_dir = os.path.join(app.root_path, 'static', 'uploads')
        os.makedirs(upload_dir, exist_ok=True)
        safe_name = secure_filename(f.filename or 'grades.xlsx')
        target_path = os.path.join(upload_dir, safe_name)
        f.save(target_path)

        students, report = _parse_summary_excel(target_path)
        mappings = _load_career_dataset()

        # Create an upload record to associate results
        meta = {'module': 'student_risk', 'source_file': safe_name}
        actor = getattr(g, 'current_user', None)
        faculty_name = None
        if actor:
            faculty_name = actor.full_name or actor.username or actor.email
        upload = FacultyUpload(faculty_name=(faculty_name or 'Unknown'), file_paths=json.dumps({'metadata': meta}))
        db.session.add(upload)
        db.session.commit()

        dist = {'Low Risk': 0, 'Medium Risk': 0, 'High Risk': 0, 'Critical': 0, 'Data Pending': 0}
        outstanding = []
        profiles = []

        for s in students:
            nm_raw = s.get('name')
            if not _is_valid_name(nm_raw):
                continue
            nm = _clean_name(nm_raw)
            breakdown = s['breakdown']
            present_grades = []
            try:
                for item in breakdown:
                    gv0 = _safe_float(item.get('grade'))
                    if gv0 is not None:
                        present_grades.append(gv0)
            except Exception:
                present_grades = []
            sanitized = []
            for item in breakdown:
                gnum = _safe_float(item.get('grade'))
                if gnum is None:
                    gnum = 0.0
                st = str(item.get('status') or '').strip().upper()
                if not st:
                    st = 'PASSED' if gnum >= 75 else 'FAILED'
                sanitized.append({'subject': item.get('subject') or '', 'grade': gnum, 'status': st})
            breakdown = sanitized
            grades_only = [b['grade'] for b in breakdown if _safe_float(b['grade']) is not None]
            avg = round(sum(grades_only)/len(grades_only), 2) if grades_only else 0.0
            try:
                import math
                if math.isnan(float(avg)):
                    avg = 0.0
            except Exception:
                avg = 0.0
            passed = sum(1 for b in breakdown if str(b['status']).upper() == 'PASSED')
            failed = sum(1 for b in breakdown if str(b['status']).upper() in ('FAILED'))
            incomplete = sum(1 for b in breakdown if str(b['status']).upper() in ('INC', 'INCOMPLETE', 'W', 'WITHDRAWN'))
            low = min(breakdown, key=lambda x: x['grade']) if breakdown else None
            high = max(breakdown, key=lambda x: x['grade']) if breakdown else None
            variance = 0.0
            if grades_only:
                mu = sum(grades_only)/len(grades_only)
                variance = sum((g-mu)**2 for g in grades_only)/len(grades_only)
            consistency = round(100.0/(1.0+variance), 2)
            all_zero = (grades_only and all(float(g or 0) == 0 for g in grades_only))
            risk = 'Data Pending' if ((len(present_grades) == 0) or all_zero) else _assign_risk(avg, failed, incomplete)
            dist[risk] = dist.get(risk, 0) + 1
            weakest_list = [low['subject']] if low else []
            # Build detailed reasoning
            triggers = [b for b in breakdown if (str(b['status']).upper() in ('FAILED','INC','INCOMPLETE','W','WITHDRAWN') or float(b['grade'] or 0) < 75)]
            trigger_str = ", ".join([f"{t['subject']}={t['grade']} ({t['status']})" for t in triggers[:4]])
            insight = _career_insight_for_student(breakdown, mappings)
            strong_subj = ", ".join((insight.get('subjects') or [])[:3])
            rec = (
                f"{nm} — {risk}. Reason: avg={avg}; failed={failed}; incomplete={incomplete}. "
                + (f"Triggers: {trigger_str}. " if trigger_str else "")
                + (f"Career: paths derived from {strong_subj}." if strong_subj else "")
            )
            profile = StudentRiskAssessment(
                upload_id=upload.id,
                student_name=nm,
                average_grade=avg,
                passed_count=passed,
                failed_count=failed,
                incomplete_count=incomplete,
                lowest_subject=(low['subject'] if low else None),
                lowest_grade=(low['grade'] if low else None),
                highest_subject=(high['subject'] if high else None),
                highest_grade=(high['grade'] if high else None),
                grade_variance=variance,
                consistency_score=consistency,
                risk_level=risk,
                recommendation=rec,
                breakdown_json=json.dumps(breakdown, allow_nan=False),
            )
            db.session.add(profile)

            # Outstanding detection
            if avg >= 90 and failed == 0 and incomplete == 0 and consistency >= 70:
                insight = _career_insight_for_student(breakdown, mappings)
                outstanding.append({
                    'name': nm,
                    'average_grade': avg,
                    'exceptional_subjects': [b['subject'] for b in breakdown if b['grade'] >= 95],
                    'remark': 'Consistently high performer across all major CS/IT competencies.',
                    'skills': insight['strong_skill_clusters'],
                })
            profiles.append(profile)

        db.session.commit()

        # Store summary analytics
        summary_payload = {
            'upload_id': upload.id,
            'distribution': dist,
            'outstanding': outstanding,
        }
        db.session.add(AnalyticsData(upload_id=upload.id, data_type='student_risk_summary', data_content=json.dumps(summary_payload)))
        # Store extraction report (audit log)
        report_payload = dict(report or {})
        report_payload['upload_id'] = upload.id
        db.session.add(AnalyticsData(upload_id=upload.id, data_type='student_risk_extraction_report', data_content=json.dumps(report_payload)))
        # Store per-student raw rows (with raw/normalized values)
        try:
            rows_payload = {'upload_id': upload.id, 'rows': students}
            db.session.add(AnalyticsData(upload_id=upload.id, data_type='student_risk_rows', data_content=json.dumps(rows_payload, default=str)))
        except Exception:
            pass
        # Store per-student career insights
        insights = {}
        for p in profiles:
            breakdown = json.loads(p.breakdown_json or '[]')
            insights[p.student_name] = _career_insight_for_student(breakdown, mappings)
        db.session.add(AnalyticsData(upload_id=upload.id, data_type='career_insight', data_content=json.dumps(insights)))
        db.session.commit()

        return jsonify({
            'upload_id': upload.id,
            'distribution': dist,
            'count': len(profiles),
        }), 201
    except Exception as e:
        db.session.rollback()
        return jsonify({'error': str(e)}), 500

@app.route('/api/student-risk/rows/latest', methods=['GET'])
def student_risk_rows_latest():
    try:
        upload = FacultyUpload.query.filter(FacultyUpload.file_paths.contains('student_risk')).order_by(FacultyUpload.analysis_date.desc()).first()
        if not upload:
            return jsonify({'rows': []})
        item = AnalyticsData.query.filter_by(upload_id=upload.id, data_type='student_risk_rows').order_by(AnalyticsData.analysis_date.desc()).first()
        payload = json.loads(item.data_content) if item else {}
        rows = payload.get('rows') or []
        return jsonify({'rows': rows, 'upload_id': payload.get('upload_id', upload.id)})
    except Exception as e:
        return jsonify({'error': str(e)}), 500

@app.route('/api/reports/list', methods=['GET'])
def list_reports():
    reports_dir = os.path.join(app.root_path, 'reports')
    os.makedirs(reports_dir, exist_ok=True)
    files = []
    try:
        for entry in os.scandir(reports_dir):
            if entry.is_file() and (entry.name.endswith('.pdf') or entry.name.endswith('.xlsx')):
                stat = entry.stat()
                files.append({
                    'filename': entry.name,
                    'type': 'Formal Report' if entry.name.endswith('.pdf') else 'Excel Export',
                    'created_at': datetime.utcfromtimestamp(stat.st_mtime).isoformat() + 'Z',
                    'size': stat.st_size
                })
    except Exception:
        pass
    files.sort(key=lambda x: x['created_at'], reverse=True)
    return jsonify(files)



@app.route('/api/student-risk/latest', methods=['GET'])
def list_student_risk_latest():
    try:
        upload = FacultyUpload.query.filter(FacultyUpload.file_paths.contains('student_risk')).order_by(FacultyUpload.analysis_date.desc()).first()
        if not upload:
            return jsonify({'items': []})
        items = StudentRiskAssessment.query.filter_by(upload_id=upload.id).all()
        return jsonify({'items': [i.to_dict() for i in items], 'upload_id': upload.id})
    except Exception as e:
        return jsonify({'error': str(e)}), 500

@app.route('/api/student-risk/summary', methods=['GET'])
def student_risk_summary():
    try:
        upload = FacultyUpload.query.filter(FacultyUpload.file_paths.contains('student_risk')).order_by(FacultyUpload.analysis_date.desc()).first()
        if not upload:
            return jsonify({'distribution': {}, 'outstanding': []})
        item = AnalyticsData.query.filter_by(upload_id=upload.id, data_type='student_risk_summary').order_by(AnalyticsData.analysis_date.desc()).first()
        payload = json.loads(item.data_content) if item else {}
        return jsonify(payload)
    except Exception as e:
        return jsonify({'error': str(e)}), 500

@app.route('/api/student-risk/leaderboard', methods=['GET'])
def student_risk_leaderboard():
    try:
        upload = FacultyUpload.query.filter(FacultyUpload.file_paths.contains('student_risk')).order_by(FacultyUpload.analysis_date.desc()).first()
        if not upload:
            return jsonify({'items': []})
        items = StudentRiskAssessment.query.filter_by(upload_id=upload.id).order_by(StudentRiskAssessment.average_grade.desc()).limit(50).all()
        ranked = []
        for i in items:
            br = []
            try:
                br = json.loads(i.breakdown_json or '[]')
            except Exception:
                br = []
            exceptional = [b['subject'] for b in br if (_safe_float(b.get('grade')) or 0) >= 95]
            ranked.append({
                'name': i.student_name,
                'average_grade': i.average_grade,
                'risk_level': i.risk_level,
                'exceptional_subjects': exceptional,
            })
        return jsonify({'items': ranked})
    except Exception as e:
        return jsonify({'error': str(e)}), 500

@app.route('/api/student-risk/insights', methods=['GET'])
def student_risk_insights():
    try:
        upload = FacultyUpload.query.filter(FacultyUpload.file_paths.contains('student_risk')).order_by(FacultyUpload.analysis_date.desc()).first()
        if not upload:
            return jsonify({'insights': {}})
        item = AnalyticsData.query.filter_by(upload_id=upload.id, data_type='career_insight').order_by(AnalyticsData.analysis_date.desc()).first()
        if not item:
            return jsonify({'insights': {}})
        return jsonify({'insights': json.loads(item.data_content)})
    except Exception as e:
        return jsonify({'error': str(e)}), 500

# Helper: Build a summary payload for a given upload_id
def compute_upload_summary(upload_id: int):
    subjects = SubjectAnalysis.query.filter_by(upload_id=upload_id).all()
    metrics_item = AnalyticsData.query.filter_by(upload_id=upload_id, data_type='precision_recall_f1').order_by(AnalyticsData.analysis_date.desc()).first()
    total_enrolled = sum(int(s.enrolled or 0) for s in subjects)
    total_passed = sum(int(s.passed or 0) for s in subjects)
    total_failed = sum(int(s.failed or 0) for s in subjects)
    overall_pass_rate = (total_passed / total_enrolled * 100.0) if total_enrolled else 0.0
    total_deficiencies = sum(int(getattr(s, 'num_def', 0) or 0) for s in subjects)
    metrics_data = {}
    if metrics_item:
        try:
            metrics_data = json.loads(metrics_item.data_content)
        except json.JSONDecodeError:
            metrics_data = {}
    upload = FacultyUpload.query.get(upload_id)
    return {
        'upload_id': upload_id,
        'faculty_name': upload.faculty_name if upload else 'Unknown',
        'analysis_date': to_manila_iso(upload.analysis_date) if upload else to_manila_iso(datetime.utcnow()),
        'subjects_processed': len(subjects),
        'total_enrolled': total_enrolled,
        'total_passed': total_passed,
        'total_failed': total_failed,
        'overall_pass_rate': round(overall_pass_rate, 2),
        'total_deficiencies': total_deficiencies,
        'metrics': metrics_data,
    }

# --- Risk Tracking module ---
@app.route('/api/risk-tracking/upload', methods=['POST'])
@roles_required('Admin', 'Faculty')
def risk_tracking_upload():
    try:
        if 'file' not in request.files:
            return jsonify({'error': 'Missing file field: file'}), 400
        f = request.files['file']
        fname = f.filename or ''
        low = fname.lower()
        if not (low.endswith('.xlsx') or low.endswith('.xls')):
            return jsonify({'error': 'Only Excel files (.xlsx, .xls) are accepted'}), 400
        upload_dir = os.path.join(app.root_path, 'static', 'uploads')
        os.makedirs(upload_dir, exist_ok=True)
        safe_name = secure_filename(fname)
        target_path = os.path.join(upload_dir, safe_name)
        f.save(target_path)

        rows = _parse_summary_for_risk_tracking(target_path)
        subject = (request.form.get('subject') or (rows[0].get('subject') if isinstance(rows, list) and len(rows) > 0 else (_extract_course_code(target_path) or 'COURSE')))
        course = request.form.get('course') or subject
        year_level = (request.form.get('year_level') or (rows[0].get('year_level') if isinstance(rows, list) and len(rows) > 0 else (_extract_year_level(target_path, None))))
        semester = request.form.get('semester') or None
        section = request.form.get('section') or None
        academic_year = request.form.get('academic_year')
        actor = getattr(g, 'current_user', None)
        faculty_name = (request.form.get('faculty_name') or (actor.full_name if actor and getattr(actor, 'full_name', None) else (actor.username if actor and getattr(actor, 'username', None) else (actor.email if actor and getattr(actor, 'email', None) else None))))

        uploader = actor
        up = StudentRiskUpload(
            uploader_id=(uploader.id if uploader else None),
            file_name=safe_name,
            subject=subject,
            course=course,
            year_level=year_level,
            semester=semester,
            section=section,
            faculty_name=faculty_name,
            academic_year=academic_year,
        )
        db.session.add(up)
        db.session.commit()

        created = 0
        for r in rows:
            grade = _safe_float(r.get('grade'))
            if grade is None:
                grade = 0.0
            st_raw = str(r.get('status') or '').strip().upper()
            status = st_raw if st_raw else ('PASSED' if float(grade) >= 75 else 'FAILED')
            rec = StudentSubjectRecord(
                upload_id=up.id,
                student_id=None,
                student_name=r.get('student_name') or '',
                subject=subject,
                year_level=(r.get('year_level') or year_level),
                grade=grade,
                status=status,
            )
            db.session.add(rec)
            created += 1

            prof = StudentAcademicProfile.query.filter_by(student_name=rec.student_name).first()
            if not prof:
                prof = StudentAcademicProfile(
                    student_id=None,
                    student_name=rec.student_name,
                    year_level=rec.year_level,
                    total_subjects=0,
                    total_grades=0.0,
                    average_grade=0.0,
                    failed_count=0,
                    passed_count=0,
                )
            prof.total_subjects = int(prof.total_subjects or 0) + 1
            prof.total_grades = float(prof.total_grades or 0.0) + float(rec.grade or 0.0)
            prof.average_grade = round((prof.total_grades / prof.total_subjects) if prof.total_subjects else 0.0, 2)
            prof.failed_count = int(prof.failed_count or 0) + (1 if status in ('FAILED','INC','INCOMPLETE','W','WITHDRAWN') else 0)
            prof.passed_count = int(prof.passed_count or 0) + (1 if status == 'PASSED' else 0)
            prof.risk_level = _assign_risk(prof.average_grade, prof.failed_count, 0)
            db.session.add(prof)

        db.session.commit()
        return jsonify({'upload_id': up.id, 'subject': subject, 'year_level': year_level, 'records_created': created}), 201
    except Exception as e:
        db.session.rollback()
        return jsonify({'error': str(e)}), 500

@app.route('/api/risk-tracking/overview', methods=['GET'])
@roles_required('Admin', 'Faculty')
def risk_tracking_overview():
    try:
        acad_year = request.args.get('academic_year')

        # Base queries
        upload_query = StudentRiskUpload.query
        if acad_year and acad_year != 'All':
            upload_query = upload_query.filter_by(academic_year=acad_year)
            
        uploads = [u.to_dict() for u in upload_query.order_by(StudentRiskUpload.uploaded_at.desc()).limit(50).all()]
        
        # For Subject Performance, filter by upload's academic year
        perf = {}
        # Join with Upload to filter
        recs_query = db.session.query(StudentSubjectRecord, StudentRiskUpload).join(StudentRiskUpload, StudentSubjectRecord.upload_id == StudentRiskUpload.id)
        if acad_year and acad_year != 'All':
            recs_query = recs_query.filter(StudentRiskUpload.academic_year == acad_year)
            
        recs = recs_query.all()
        for r, u in recs:
            key = r.subject
            d = perf.get(key) or {'subject': key, 'count': 0, 'passed': 0, 'failed': 0}
            d['count'] += 1
            if str(r.status).upper() == 'PASSED':
                d['passed'] += 1
            else:
                d['failed'] += 1
            perf[key] = d

        # Get available years
        all_years = [r.academic_year for r in db.session.query(StudentRiskUpload.academic_year).distinct().all() if r.academic_year]

        profiles = StudentAcademicProfile.query.all()
        dist = {'Low Risk': 0, 'Medium Risk': 0, 'High Risk': 0, 'Critical': 0}
        for p in profiles:
            if p.risk_level:
                dist[p.risk_level] = dist.get(p.risk_level, 0) + 1
        outstanding = [p.to_dict() for p in profiles if (float(p.average_grade or 0) >= 90 and int(p.failed_count or 0) == 0)]
        
        return jsonify({
            'risk_distribution': dist, 
            'outstanding_count': len(outstanding), 
            'uploads': uploads, 
            'subject_performance': list(perf.values()),
            'years': sorted(all_years, reverse=True)
        })
    except Exception as e:
        return jsonify({'error': str(e)}), 500

@app.route('/api/risk-tracking/top5', methods=['GET'])
@roles_required('Admin', 'Faculty')
def risk_tracking_top5():
    try:
        subj = (request.args.get('subject') or '').strip()
        def rank_for_subject(s):
            q = StudentSubjectRecord.query.filter_by(subject=s)
            items = [r.to_dict() for r in q.all()]
            items = [
                {'student_name': it['student_name'], 'grade': float(it['grade'] or 0.0)}
                for it in items
                if it.get('student_name') and (it.get('grade') is not None)
            ]
            ranked = sorted(items, key=lambda x: (-float(x['grade']), x['student_name']))[:5]
            out = []
            for idx, it in enumerate(ranked, start=1):
                out.append({'rank': idx, 'student_name': it['student_name'], 'grade': float(it['grade'])})
            return out
        if subj:
            return jsonify({'subject': subj, 'top5': rank_for_subject(subj)})
        subs = [s.subject for s in db.session.query(StudentSubjectRecord.subject).distinct().all()]
        items = []
        for s in subs:
            items.append({'subject': s, 'top5': rank_for_subject(s)})
        return jsonify({'items': items})
    except Exception as e:
        return jsonify({'error': str(e)}), 500

@app.route('/api/risk-tracking/subjects', methods=['GET'])
@roles_required('Admin', 'Faculty')
def risk_tracking_subjects():
    try:
        from sqlalchemy import func
        rows = db.session.query(StudentSubjectRecord.subject, func.count(StudentSubjectRecord.id)).group_by(StudentSubjectRecord.subject).all()
        return jsonify({'items': [{'subject': s, 'count': int(c)} for s, c in rows]})
    except Exception as e:
        return jsonify({'error': str(e)}), 500

@app.route('/api/risk-tracking/subject-records', methods=['GET'])
@roles_required('Admin', 'Faculty')
def risk_tracking_subject_records():
    try:
        subj = (request.args.get('subject') or '').strip()
        q = StudentSubjectRecord.query
        if subj:
            q = q.filter_by(subject=subj)
        items = [r.to_dict() for r in q.order_by(StudentSubjectRecord.created_at.desc()).limit(1000).all()]
        return jsonify({'items': items})
    except Exception as e:
        return jsonify({'error': str(e)}), 500

@app.route('/api/risk-tracking/outstanding', methods=['GET'])
@roles_required('Admin', 'Faculty')
def risk_tracking_outstanding():
    try:
        yl = (request.args.get('year') or '').strip()
        q = StudentAcademicProfile.query
        if yl:
            q = q.filter_by(year_level=yl)
        items = [p.to_dict() for p in q.all() if (float(p.average_grade or 0) >= 90 and int(p.failed_count or 0) == 0)]
        items = sorted(items, key=lambda x: (-float(x['average_grade'] or 0), x['student_name']))
        return jsonify({'items': items})
    except Exception as e:
        return jsonify({'error': str(e)}), 500

@app.route('/api/risk-tracking/profiles', methods=['GET'])
@roles_required('Admin', 'Faculty')
def risk_tracking_profiles():
    try:
        items = [p.to_dict() for p in StudentAcademicProfile.query.order_by(StudentAcademicProfile.average_grade.desc()).all()]
        return jsonify({'items': items})
    except Exception as e:
        return jsonify({'error': str(e)}), 500

@app.route('/api/risk-tracking/profile/<name>', methods=['GET'])
@roles_required('Admin', 'Faculty')
def risk_tracking_profile(name):
    try:
        prof = StudentAcademicProfile.query.filter_by(student_name=name).first()
        records = [r.to_dict() for r in StudentSubjectRecord.query.filter_by(student_name=name).order_by(StudentSubjectRecord.created_at.desc()).all()]
        br = [{'subject': r['subject'], 'grade': r['grade'], 'status': r['status']} for r in records]
        mappings = _load_career_dataset()
        insight = _career_insight_for_student(br, mappings)
        profile_payload = (prof.to_dict() if prof else {'student_name': name, 'average_grade': None, 'passed_count': None, 'failed_count': None, 'risk_level': 'Low Risk'})
        try:
            avg = profile_payload.get('average_grade')
            avg_str = f"{float(avg):.2f}" if avg is not None else '—'
            risk = profile_payload.get('risk_level') or '—'
            prefix = f"Risk level: {risk}; Avg grade: {avg_str}"
            prev = insight.get('reasoning') or ''
            insight['reasoning'] = f"{prefix}; {prev}" if prev else prefix
        except Exception:
            pass
        return jsonify({'profile': profile_payload, 'records': records, 'career_insight': insight})
    except Exception as e:
        return jsonify({'error': str(e)}), 500

@app.route('/api/risk-tracking/profiles-insights', methods=['POST'])
@roles_required('Admin', 'Faculty')
def risk_tracking_profiles_insights():
    try:
        payload = request.get_json(force=True) or {}
        names = payload.get('names') or []
        names = [str(n).strip() for n in names if str(n or '').strip()]
        if not names:
            return jsonify({'items': []})
        q = StudentSubjectRecord.query.filter(StudentSubjectRecord.student_name.in_(names))
        recs = [r.to_dict() for r in q.all()]
        by_name = {}
        for r in recs:
            nm = r.get('student_name')
            arr = by_name.get(nm) or []
            arr.append({'subject': r.get('subject'), 'grade': r.get('grade'), 'status': r.get('status')})
            by_name[nm] = arr
        mappings = _load_career_dataset()
        items = []
        for n in names:
            br = by_name.get(n) or []
            insight = _career_insight_for_student(br, mappings)
            items.append({'student_name': n, 'career_insight': insight})
        return jsonify({'items': items})
    except Exception as e:
        return jsonify({'error': str(e)}), 500

@app.route('/')
def index():
    return redirect(url_for('static', filename='index.html')) # Redirect root to frontend if served separately

# --- Auth API ---
@app.route('/api/auth/login', methods=['POST'])
def login():
    payload = request.get_json(force=True) or {}
    identifier = (payload.get('identifier') or '').strip()
    password = (payload.get('password') or '').strip()
    selected_role = (payload.get('role') or '').strip()
    if not identifier or not password:
        return jsonify({'error': 'Missing credentials'}), 400

    user = None
    if '@' in identifier:
        user = User.query.filter_by(email=identifier).first()
    else:
        user = User.query.filter_by(username=identifier).first()
    if not user or not user.active:
        return jsonify({'error': 'Invalid credentials'}), 401
    if selected_role and user.role != selected_role:
        return jsonify({'error': 'Role mismatch'}), 403
    if not check_password_hash(user.password_hash, password):
        return jsonify({'error': 'Invalid credentials'}), 401
    token = issue_token(user)
    return jsonify({'token': token, 'user': user.to_public()})

# Registration for Faculty and Student
@app.route('/api/auth/register', methods=['POST'])
def register():
    payload = request.get_json(force=True) or {}
    role = (payload.get('role') or '').strip()
    email = (payload.get('email') or '').strip().lower()
    password = (payload.get('password') or '')
    confirm = (payload.get('confirmPassword') or '')

    # Common validations
    if role not in ('Faculty', 'Student'):
        return jsonify({'error': 'Invalid role for self-registration'}), 400
    if not email or not password or not confirm:
        return jsonify({'error': 'email, password, confirmPassword required'}), 400
    if password != confirm:
        return jsonify({'error': 'Password and confirmation do not match'}), 400
    # Minimal password strength: 8+, upper, lower, digit
    has_upper = any(c.isupper() for c in password)
    has_lower = any(c.islower() for c in password)
    has_digit = any(c.isdigit() for c in password)
    if len(password) < 8 or not (has_upper and has_lower and has_digit):
        return jsonify({'error': 'Password must be at least 8 chars and include upper, lower, digit'}), 400

    # Role-specific fields
    full_name = (payload.get('fullName') or '').strip()
    faculty_id = (payload.get('facultyId') or '').strip() if role == 'Faculty' else None
    student_id = (payload.get('studentId') or '').strip() if role == 'Student' else None
    program = (payload.get('program') or '').strip() if role == 'Student' else None

    if not full_name:
        return jsonify({'error': 'Full Name is required'}), 400
    if role == 'Faculty' and not faculty_id:
        return jsonify({'error': 'Faculty ID is required'}), 400
    if role == 'Student' and (not student_id or not program):
        return jsonify({'error': 'Student ID and Program are required'}), 400

    # Uniqueness checks
    if User.query.filter((User.email == email)).first():
        return jsonify({'error': 'Email already exists'}), 409
    if faculty_id and User.query.filter(User.faculty_id == faculty_id).first():
        return jsonify({'error': 'Faculty ID already exists'}), 409
    if student_id and User.query.filter(User.student_id == student_id).first():
        return jsonify({'error': 'Student ID already exists'}), 409

    # Use email as username for simplicity
    username = email
    if User.query.filter(User.username == username).first():
        return jsonify({'error': 'Username already exists'}), 409

    user = User(
        username=username,
        email=email,
        role=role,
        password_hash=generate_password_hash(password),
        full_name=full_name,
        faculty_id=faculty_id,
        student_id=student_id,
        program=program,
        active=True,
    )
    db.session.add(user)
    db.session.commit()
    # Return minimal user and instruct client to login
    return jsonify({'user': user.to_public()}), 201

@app.route('/api/auth/me', methods=['GET'])
@auth_required
def me():
    user = g.current_user
    return jsonify({'user': user.to_public()})

@app.route('/api/auth/logout', methods=['POST'])
def logout():
    # Client-side should discard the token; provided for symmetry
    return jsonify({'ok': True})

# Admin: create or manage users
@app.route('/api/admin/users', methods=['POST'])
@roles_required('Admin')
def create_user():
    payload = request.get_json(force=True) or {}
    username = (payload.get('username') or '').strip()
    email = (payload.get('email') or '').strip()
    role = (payload.get('role') or '').strip()
    password = (payload.get('password') or '').strip()
    if not username or not email or not role or not password:
        return jsonify({'error': 'username, email, role, password required'}), 400
    if role not in ('Admin', 'Faculty', 'Student'):
        return jsonify({'error': 'Invalid role'}), 400
    if User.query.filter((User.username == username) | (User.email == email)).first():
        return jsonify({'error': 'User already exists'}), 409
    user = User(
        username=username,
        email=email,
        role=role,
        password_hash=generate_password_hash(password)
    )
    db.session.add(user)
    db.session.commit()
    return jsonify({'user': user.to_public()}), 201

# API Endpoint: Get summary data for the home page
@app.route('/api/dashboard', methods=['GET'])
def get_dashboard_data():
    uploads = FacultyUpload.query.order_by(FacultyUpload.analysis_date.desc()).all()
    # Deduplicate recent uploads by idempotency key scoped to faculty (keep newest)
    deduped_uploads = []
    seen = set()
    for u in uploads:
        key = None
        try:
            meta = json.loads(u.file_paths or '{}')
            idem = (meta.get('metadata', {}) or {}).get('idempotency_key') or ''
            if idem:
                key = f"{u.faculty_name}:{idem}"
        except Exception:
            key = None
        if not key:
            key = f"{u.faculty_name}:{to_manila_iso(u.analysis_date)}"
        if key in seen:
            continue
        seen.add(key)
        deduped_uploads.append(u)
    subjects = SubjectAnalysis.query.all()

    total_uploads = len(deduped_uploads)
    total_subjects = len(subjects)
    total_enrolled = sum([s.enrolled for s in subjects])
    total_passed = sum([s.passed for s in subjects])
    total_failed = sum([s.failed for s in subjects])
    total_deficiencies = sum([s.num_def for s in subjects])
    overall_pass_rate = (total_passed / total_enrolled * 100) if total_enrolled > 0 else 0

    # Fetch latest metrics
    latest_metrics = AnalyticsData.query.filter_by(data_type='precision_recall_f1').order_by(AnalyticsData.analysis_date.desc()).first()
    metrics = {}
    if latest_metrics:
        try:
            metrics = json.loads(latest_metrics.data_content)
        except json.JSONDecodeError:
            metrics = {}

    return jsonify({
        'total_uploads': total_uploads,
        'total_subjects': total_subjects,
        'total_enrolled': total_enrolled,
        'total_passed': total_passed,
        'total_failed': total_failed,
        'total_deficiencies': total_deficiencies,
        'overall_pass_rate': round(overall_pass_rate, 2),
        'recent_uploads': [upload.to_dict() for upload in deduped_uploads[:50]],
        'latest_metrics': metrics
    })

# NEW API Endpoint: List generated reports
@app.route('/api/reports/list', methods=['GET'])
def list_generated_reports():
    reports_dir = os.path.join(app.root_path, 'static', 'reports')
    os.makedirs(reports_dir, exist_ok=True)
    
    # Also check results directory for AI generated reports
    results_dir = os.path.join(app.root_path, 'results')
    
    reports = []
    
    # List manually generated reports
    for f in os.listdir(reports_dir):
        if f.startswith('.'): continue
        path = os.path.join(reports_dir, f)
        if os.path.isfile(path):
            stat = os.stat(path)
            reports.append({
                'filename': f,
                'path': f'/static/reports/{f}',
                'type': 'Generated',
                'created_at': datetime.fromtimestamp(stat.st_mtime).isoformat(),
                'size': stat.st_size
            })
            
    # List AI processor results (recursively if needed, but keeping it simple)
    if os.path.exists(results_dir):
        for root, dirs, files in os.walk(results_dir):
            for f in files:
                if f.endswith('.csv') or f.endswith('.xlsx'):
                    path = os.path.join(root, f)
                    rel_path = os.path.relpath(path, app.root_path) # e.g. results/Name/file.csv
                    # We need a route to serve these if they are not in static
                    # Or we can copy them to static/reports on demand?
                    # For now, let's assume we can serve them via a specific endpoint or just ignore if hard to serve
                    # Actually, let's just list them and maybe add a download route for results
                    stat = os.stat(path)
                    reports.append({
                        'filename': f,
                        'path': f'/api/reports/download_result?path={requests.utils.quote(os.path.relpath(path, results_dir))}', # We need a helper for this
                        'type': 'AI Result',
                        'created_at': datetime.fromtimestamp(stat.st_mtime).isoformat(),
                        'size': stat.st_size
                    })

    # Sort by created_at desc
    reports.sort(key=lambda x: x['created_at'], reverse=True)
    return jsonify(reports)

# NEW API Endpoint: Download result file
@app.route('/api/reports/download_result', methods=['GET'])
def download_result_file():
    rel_path = request.args.get('path', '')
    if not rel_path or '..' in rel_path:
        return jsonify({'error': 'Invalid path'}), 400
    
    results_dir = os.path.join(app.root_path, 'results')
    safe_path = os.path.join(results_dir, rel_path)
    
    if not os.path.exists(safe_path) or not os.path.isfile(safe_path):
        return jsonify({'error': 'File not found'}), 404
        
    return send_file(safe_path, as_attachment=True)



# NEW API Endpoint: Generate Report
@app.route('/api/reports/generate', methods=['POST'])
def generate_report():
    import pandas as pd
    data = request.json or {}
    report_type = data.get('type', 'performance_summary')
    academic_year = data.get('year', 'All')
    
    # Prepare data
    subjects_query = SubjectAnalysis.query
    if academic_year != 'All':
        # Filter by year logic (assuming upload has analysis_date or similar)
        # This is tricky without explicit year column, but we can filter by upload date roughly
        # Or if we had academic_year in Upload
        pass
        
    subjects = subjects_query.all()
    
    timestamp = datetime.now().strftime('%Y%m%d_%H%M%S')
    filename = f"Report_{report_type}_{timestamp}.xlsx"
    filepath = os.path.join(app.root_path, 'static', 'reports', filename)
    
    if report_type == 'performance_summary':
        # Create summary DataFrame
        records = []
        for s in subjects:
            records.append({
                'Course': s.course,
                'Program': s.program,
                'Enrolled': s.enrolled,
                'Passed': s.passed,
                'Failed': s.failed,
                'Pass Rate': s.pass_rate,
                'Deficiencies': s.num_def,
                'Recommendation': s.recommendation
            })
        df = pd.DataFrame(records)
        
        # Calculate summary stats
        summary_stats = pd.DataFrame([{
            'Total Students': df['Enrolled'].sum(),
            'Total Passed': df['Passed'].sum(),
            'Total Failed': df['Failed'].sum(),
            'Average Pass Rate': df['Pass Rate'].mean()
        }])
        
        # Save to Excel with multiple sheets
        with pd.ExcelWriter(filepath) as writer:
            summary_stats.to_excel(writer, sheet_name='Summary', index=False)
            df.to_excel(writer, sheet_name='Detailed Data', index=False)
            
    elif report_type == 'faculty_evaluation':
        # Group by faculty
        records = []
        for s in subjects:
            faculty = s.upload.faculty_name if s.upload else 'Unknown'
            records.append({
                'Faculty': faculty,
                'Course': s.course,
                'Pass Rate': s.pass_rate,
                'Enrolled': s.enrolled
            })
        df = pd.DataFrame(records)
        if not df.empty:
            summary = df.groupby('Faculty').agg({
                'Course': 'count',
                'Enrolled': 'sum',
                'Pass Rate': 'mean'
            }).reset_index()
            summary.columns = ['Faculty', 'Subjects Taught', 'Total Students', 'Avg Pass Rate']
            
            with pd.ExcelWriter(filepath) as writer:
                summary.to_excel(writer, sheet_name='Faculty Summary', index=False)
                df.to_excel(writer, sheet_name='Raw Data', index=False)
        else:
             pd.DataFrame({'Message': ['No data available']}).to_excel(filepath, index=False)

    elif report_type == 'student_risk':
        # Fetch student risk data
        # This might need parsing stored JSONs
        pass
        
    return jsonify({
        'message': 'Report generated successfully',
        'filename': filename,
        'url': f'/static/reports/{filename}'
    })

# API Endpoint: Get all subjects and their recommendations
@app.route('/api/subjects', methods=['GET'])
def get_subjects():
    from models import SubjectAccomplishment
    subjects = SubjectAnalysis.query.all()
    
    # Fetch all accomplishment records
    acc_records = SubjectAccomplishment.query.all()
    acc_map = {}
    for r in acc_records:
        if not r.subject_code: continue
        # Normalize subject code for matching
        k = r.subject_code.upper().replace(' ', '')
        if k not in acc_map: acc_map[k] = []
        acc_map[k].append(r.to_dict())
    
    results = []
    for subject in subjects:
        data = subject.to_dict()
        # Find matching accomplishment
        # Try exact match first, then fuzzy or containment
        s_code = subject.course.upper().replace(' ', '')
        
        # Simple containment matching
        matched_reviews = []
        for code_key, reviews in acc_map.items():
            if code_key in s_code or s_code in code_key:
                matched_reviews.extend(reviews)
        
        # Deduplicate by id if needed, but for now just pass list
        # Sort by latest created_at
        matched_reviews.sort(key=lambda x: x['created_at'], reverse=True)
        
        data['internal_reviews'] = matched_reviews
        results.append(data)
        
    return jsonify(results)

# API Endpoint: Get analytics data
@app.route('/api/analysis', methods=['GET'])
def get_analysis():
    subjects = SubjectAnalysis.query.all()
    analytics_data = AnalyticsData.query.all()

    # Process stored JSON data if needed
    processed_analytics = []
    for item in analytics_data:
        try:
            content = json.loads(item.data_content)
        except json.JSONDecodeError:
            content = item.data_content # If not JSON, keep as string
        processed_analytics.append({
            'id': item.id,
            'upload_id': item.upload_id,
            'data_type': item.data_type,
            'content': content,
            'analysis_date': item.analysis_date.isoformat()
        })

    # Prepare subject data for charts
    subject_chart_data = [
        {
            'name': s.course,
            'passRate': s.pass_rate,
            'enrolled': s.enrolled,
            'deficiencies': getattr(s, 'num_def', 0),
            'failed': s.failed,
            'passed': s.passed
        }
        for s in subjects
    ]

    # Prepare category data for charts
    category_performance = {}
    for s in subjects:
        # Simplified category logic based on course name
        cat = "General"
        if any(kw in s.course.upper() for kw in ['MOBILE', 'COMPUTING', 'LAB', 'PRACTICUM', 'INTEGRATION']):
            cat = "Lab"
        elif any(kw in s.course.upper() for kw in ['THEORY', 'PRINCIPLES', 'CONCEPTS', 'WORLD']):
            cat = "Theory"
        elif any(kw in s.course.upper() for kw in ['SOFTWARE', 'ENGINEERING', 'ALGORITHM', 'COMPLEXITY']):
            cat = "Software"
        elif any(kw in s.course.upper() for kw in ['APPLICATION', 'DEVELOPMENT', 'MULTIMEDIA']):
            cat = "Development"
        elif any(kw in s.course.upper() for kw in ['DATABASE', 'WEB', 'CLOUD']):
            cat = "Web/DB"
        elif any(kw in s.course.upper() for kw in ['ALGORITHM', 'DATA STRUCTURE']):
            cat = "Algorithm"
        elif any(kw in s.course.upper() for kw in ['ETHICS', 'ASSURANCE', 'SECURITY']):
            cat = "Security"
        elif any(kw in s.course.upper() for kw in ['COMMUNICATION', 'ART', 'APPRECIATION']):
            cat = "Humanities"

        if cat not in category_performance:
            category_performance[cat] = {'total_pass_rate': 0, 'count': 0}
        category_performance[cat]['total_pass_rate'] += s.pass_rate
        category_performance[cat]['count'] += 1

    category_chart_data = []
    for cat, data in category_performance.items():
        avg_pass_rate = data['total_pass_rate'] / data['count'] if data['count'] > 0 else 0
        category_chart_data.append({
            'name': cat,
            'avg_pass_rate': avg_pass_rate
        })

    return jsonify({
        'analytics': processed_analytics,
        'subjects': subject_chart_data,
        'categories': category_chart_data
    })

# NEW API Endpoint: Get the detailed analytics output string
@app.route('/api/analysis/detailed-output', methods=['GET'])
def get_detailed_analytics():
    # Fetch the most recent detailed analytics output
    # You might want to filter by upload_id if you want output for a specific upload
    latest_output = AnalyticsData.query.filter_by(data_type='detailed_analytics_output').order_by(AnalyticsData.analysis_date.desc()).first()

    if latest_output:
        return jsonify({'detailed_output': latest_output.data_content})
    else:
        return jsonify({'detailed_output': 'No detailed analytics output found.'}), 404

# Admin Endpoint: Clear analytics outputs (use with caution)
@app.route('/api/admin/clear-analytics', methods=['GET'])
@roles_required('Admin')
def clear_analytics():
    confirm = request.args.get('confirm', '').lower()
    if confirm not in ('1', 'true', 'yes'): 
        return jsonify({'error': "Confirmation required. Append '?confirm=1' or 'true'"}), 400

    types_param = request.args.get('types', '').strip()
    if types_param:
        types = [t.strip() for t in types_param.split(',') if t.strip()]
    else:
        # Default: clear detailed output and metrics
        types = ['detailed_analytics_output', 'precision_recall_f1']

    total_deleted = 0
    details = []
    for t in types:
        deleted = AnalyticsData.query.filter_by(data_type=t).delete(synchronize_session=False)
        total_deleted += deleted
        details.append({'data_type': t, 'deleted': deleted})
    db.session.commit()

    return jsonify({'deleted_total': total_deleted, 'details': details})

# Admin Endpoint: Delete uploaded/result files (use with caution)
@app.route('/api/admin/clear-files', methods=['GET'])
@roles_required('Admin')
def clear_files():
    confirm = request.args.get('confirm', '').lower()
    if confirm not in ('1', 'true', 'yes'):
        return jsonify({'error': "Confirmation required. Append '?confirm=1' or 'true'"}), 400

    scope = (request.args.get('scope', 'uploads') or 'uploads').lower()
    filename = (request.args.get('filename', '') or '').strip()

    upload_dir = os.path.join(app.root_path, 'static', 'uploads')
    results_dir = os.path.join(app.root_path, 'results')

    targets = []
    if scope == 'uploads':
        targets = [upload_dir]
    elif scope == 'results':
        targets = [results_dir]
    elif scope == 'all':
        targets = [upload_dir, results_dir]
    else:
        return jsonify({'error': "Invalid 'scope'. Use 'uploads', 'results', or 'all'."}), 400

    deleted = []
    errors = []

    def safe_delete_file(base_dir, name):
        try:
            candidate = os.path.abspath(os.path.join(base_dir, name))
            base_abs = os.path.abspath(base_dir)
            if not candidate.startswith(base_abs):
                raise ValueError('Path traversal detected')
            if os.path.isfile(candidate):
                os.remove(candidate)
                deleted.append(candidate)
                return True
            else:
                errors.append({'file': candidate, 'error': 'Not a file'})
                return False
        except Exception as e:
            errors.append({'file': os.path.join(base_dir, name), 'error': str(e)})
            return False

    def delete_all_files_in_dir(base_dir):
        try:
            if not os.path.isdir(base_dir):
                return
            for root, _, files in os.walk(base_dir):
                for f in files:
                    fp = os.path.join(root, f)
                    try:
                        os.remove(fp)
                        deleted.append(os.path.abspath(fp))
                    except Exception as e:
                        errors.append({'file': fp, 'error': str(e)})
        except Exception as e:
            errors.append({'dir': base_dir, 'error': str(e)})

    # Perform deletion
    if filename:
        for t in targets:
            safe_delete_file(t, filename)
    else:
        for t in targets:
            delete_all_files_in_dir(t)

    return jsonify({
        'scope': scope,
        'filename': filename or None,
        'deleted_count': len(deleted),
        'deleted': deleted[:50],  # cap to avoid huge responses
        'errors': errors
    })

# Admin Endpoint: List uploaded/result files
@app.route('/api/admin/list-files', methods=['GET'])
@roles_required('Admin')
def list_files():
    scope = (request.args.get('scope', 'uploads') or 'uploads').lower()

    upload_dir = os.path.join(app.root_path, 'static', 'uploads')
    results_dir = os.path.join(app.root_path, 'results')

    def list_dir_files(dir_path):
        os.makedirs(dir_path, exist_ok=True)
        files = []
        try:
            for entry in os.scandir(dir_path):
                if entry.is_file():
                    stat = entry.stat()
                    files.append({
                        'name': entry.name,
                        'size': stat.st_size,
                        'modified': datetime.utcfromtimestamp(stat.st_mtime).isoformat() + 'Z'
                    })
        except Exception as e:
            return {'error': str(e)}
        return files

    if scope == 'uploads':
        return jsonify({'scope': 'uploads', 'files': list_dir_files(upload_dir)})
    elif scope == 'results':
        return jsonify({'scope': 'results', 'files': list_dir_files(results_dir)})
    elif scope == 'all':
        return jsonify({
            'scope': 'all',
            'uploads': list_dir_files(upload_dir),
            'results': list_dir_files(results_dir)
        })
    else:
        return jsonify({'error': "Invalid 'scope'. Use 'uploads', 'results', or 'all'."}), 400

# Admin Endpoint: Full purge (DB tables + files)
@app.route('/api/admin/reset-all', methods=['GET'])
@roles_required('Admin')
def reset_all():
    confirm = request.args.get('confirm', '').lower()
    if confirm not in ('1', 'true', 'yes'):
        return jsonify({'error': "Confirmation required. Append '?confirm=1' or 'true'"}), 400

    # Purge database tables
    deleted = {
        'StudentSubjectRecord': 0,
        'StudentAcademicProfile': 0,
        'StudentRiskUpload': 0,
        'StudentRiskAssessment': 0,
        'SubjectAnalysis': 0,
        'AnalyticsData': 0,
        'UniversalUploadLog': 0,
        'FacultyUpload': 0,
        'IdempotencyKey': 0,
    }
    try:
        # Delete dependents first to satisfy foreign key constraints
        deleted['StudentSubjectRecord'] = StudentSubjectRecord.query.delete(synchronize_session=False)
        deleted['StudentAcademicProfile'] = StudentAcademicProfile.query.delete(synchronize_session=False)
        deleted['StudentRiskUpload'] = StudentRiskUpload.query.delete(synchronize_session=False)
        deleted['StudentRiskAssessment'] = StudentRiskAssessment.query.delete(synchronize_session=False)
        deleted['SubjectAnalysis'] = SubjectAnalysis.query.delete(synchronize_session=False)
        deleted['AnalyticsData'] = AnalyticsData.query.delete(synchronize_session=False)
        deleted['UniversalUploadLog'] = UniversalUploadLog.query.delete(synchronize_session=False)
        deleted['IdempotencyKey'] = IdempotencyKey.query.delete(synchronize_session=False)
        deleted['FacultyUpload'] = FacultyUpload.query.delete(synchronize_session=False)
        db.session.commit()
    except Exception as e:
        db.session.rollback()
        return jsonify({'error': f'DB purge failed: {str(e)}'}), 500

    # Purge files in uploads and results
    upload_dir = os.path.join(app.root_path, 'static', 'uploads')
    results_dir = os.path.join(app.root_path, 'results')
    file_deleted = []
    file_errors = []

    def delete_all_files_in_dir(base_dir):
        try:
            if not os.path.isdir(base_dir):
                return
            for root, _, files in os.walk(base_dir):
                for f in files:
                    fp = os.path.join(root, f)
                    try:
                        os.remove(fp)
                        file_deleted.append(os.path.abspath(fp))
                    except Exception as e:
                        file_errors.append({'file': fp, 'error': str(e)})
        except Exception as e:
            file_errors.append({'dir': base_dir, 'error': str(e)})

    delete_all_files_in_dir(upload_dir)
    delete_all_files_in_dir(results_dir)

    return jsonify({
        'db_deleted': deleted,
        'files_deleted_count': len(file_deleted),
        'files_deleted': file_deleted[:50],
        'file_errors': file_errors,
        'message': 'System data reset complete'
    })


# API Endpoint: Upload and process files
@app.route('/api/upload', methods=['POST'])
@roles_required('Admin', 'Faculty')
def upload_files():
    # Check if this is an accomplishment report upload
    if 'subject_accomplishment' in request.files:
        f = request.files['subject_accomplishment']
        upload_dir = os.path.join(app.root_path, 'static', 'uploads')
        os.makedirs(upload_dir, exist_ok=True)
        
        safe_name = secure_filename(f.filename or 'accomplishment.pdf')
        target_path = os.path.join(upload_dir, safe_name)
        f.save(target_path)
        
        # Create upload record
        faculty_name = request.form.get('faculty_name', 'Unknown')
        academic_year = request.form.get('academic_year')
        meta = {'module': 'subject_accomplishment', 'source_file': safe_name}
        upload = FacultyUpload(
            faculty_name=faculty_name, 
            file_paths=json.dumps({'metadata': meta}),
            academic_year=academic_year
        )
        db.session.add(upload)
        db.session.commit()
        
        # Process
        from ai_processor import process_subject_accomplishment
        ftype = 'pdf'
        if safe_name.lower().endswith('.docx'): ftype = 'docx'
        elif safe_name.lower().endswith('.xlsx') or safe_name.lower().endswith('.xls'): ftype = 'xlsx'
            
        process_subject_accomplishment(target_path, ftype, upload.id)
        
        return jsonify({'message': 'Accomplishment report uploaded and processed successfully', 'upload_id': upload.id})

    if 'class_profile' not in request.files or 'def_report' not in request.files:
        return jsonify({'error': 'Please upload both Class Profile and Deficiency Report files.'}), 400

    class_profile = request.files['class_profile']
    def_report = request.files['def_report']
    class_records = request.files.getlist('class_records') # Get list of files
    faculty_name = request.form.get('faculty_name', 'Unknown')
    program = request.form.get('program')
    semester = request.form.get('semester')
    academic_year = request.form.get('academic_year')
    idem_key = request.form.get('idempotency_key')
    scoped_key = None
    if idem_key:
        scoped_key = f"{faculty_name}:{idem_key}"

    # --- Save uploaded files ---
    upload_dir = os.path.join(app.root_path, 'static', 'uploads')
    os.makedirs(upload_dir, exist_ok=True)

    class_profile_path = os.path.join(upload_dir, class_profile.filename)
    def_report_path = os.path.join(upload_dir, def_report.filename)
    class_profile.save(class_profile_path)
    def_report.save(def_report_path)

    # Process Class Record files if any
    class_record_path_map = {}
    for file in class_records:
        if file.filename != '':
            record_path = os.path.join(upload_dir, file.filename)
            file.save(record_path)
            # Extract course code from filename (COURSE_CODE_Class_Record_SECTION.xlsx)
            import re
            match = re.match(r'^([A-Z0-9]+)_Class_Record_([A-Z0-9]+)\.xlsx$', file.filename, re.IGNORECASE)
            if match:
                course_code = match.group(1)
                class_record_path_map[record_path] = course_code
            else:
                print(f"  Warning: File {file.filename} does not match expected Class Record naming pattern. Skipping as Class Record.")

    # --- Idempotency handling: if key exists, remove previous records and reprocess ---
    existing_map = None
    if scoped_key:
        try:
            existing_map = IdempotencyKey.query.filter_by(key=scoped_key).first()
            if existing_map:
                try:
                    # Delete previous records tied to this idempotency mapping
                    SubjectAnalysis.query.filter_by(upload_id=existing_map.upload_id).delete(synchronize_session=False)
                    AnalyticsData.query.filter_by(upload_id=existing_map.upload_id).delete(synchronize_session=False)
                    FacultyUpload.query.filter_by(id=existing_map.upload_id).delete(synchronize_session=False)
                    db.session.commit()
                except Exception as del_err:
                    # Roll back deletion attempt and proceed to reprocess anyway
                    db.session.rollback()
        except Exception as lookup_err:
            # Non-blocking: continue without idempotency cleanup
            pass

    # --- Run the AI analysis ---
    try:
        # Pass the scoped key so metadata uniquely identifies re-uploads per faculty
        result = run_analysis(class_profile_path, def_report_path, faculty_name, class_record_path_map, program, semester, scoped_key, academic_year=academic_year)
        # Persist or update idempotency mapping for future requests
        if scoped_key:
            try:
                new_upload_id = int(result.get('upload_id'))
                if existing_map:
                    # Update existing mapping to point to the new upload
                    existing_map.upload_id = new_upload_id
                    db.session.commit()
                else:
                    db.session.add(IdempotencyKey(key=scoped_key, upload_id=new_upload_id))
                    db.session.commit()
            except Exception as id_err:
                # On failure, roll back and continue returning the fresh result
                db.session.rollback()
        return jsonify(result), 201
    except Exception as e:
        print(f"Error processing files: {e}")
        return jsonify({'error': f'Error processing files: {str(e)}'}), 500


# --- Universal Uploads: accept any file types ---
def _detect_mime_and_size(filepath: str, fallback_filename: str, incoming_mimetype: str = None):
    import mimetypes
    guessed_type, _ = mimetypes.guess_type(fallback_filename)
    mime = incoming_mimetype or guessed_type or 'application/octet-stream'
    try:
        size = os.path.getsize(filepath)
    except Exception:
        size = None
    return mime, size

def _extract_text_snippet_from_path(filepath: str, mime: str):
    snippet = None
    try:
        # Plain text-like content
        if mime.startswith('text/') or any(filepath.lower().endswith(ext) for ext in ('.txt', '.md', '.css', '.js', '.rtf', '.svg')):
            try:
                with open(filepath, 'r', encoding='utf-8', errors='ignore') as f:
                    data = f.read(2000)
                    # For RTF, try to strip control words if library is present
                    if mime == 'text/rtf' or filepath.lower().endswith('.rtf'):
                        try:
                            from striprtf.striprtf import rtf_to_text
                            data = rtf_to_text(data)
                        except Exception:
                            pass
                    snippet = data[:1200]
            except Exception:
                snippet = None

        # JSON, XML, HTML
        elif mime in ('application/json', 'application/xml', 'text/xml', 'text/html', 'application/xhtml+xml'):
            try:
                with open(filepath, 'r', encoding='utf-8', errors='ignore') as f:
                    data = f.read(2000)
                    snippet = data[:1200]
            except Exception:
                snippet = None

        # CSV / Excel
        elif mime in (
            'application/vnd.ms-excel',
            'application/vnd.openxmlformats-officedocument.spreadsheetml.sheet',
            'text/csv',
        ) or filepath.lower().endswith(('.xls', '.xlsx', '.csv')):
            import pandas as pd
            try:
                if mime == 'text/csv' or filepath.lower().endswith('.csv'):
                    df = pd.read_csv(filepath)
                else:
                    df = pd.read_excel(filepath)
                snippet = df.head(10).to_string()
            except Exception:
                snippet = None

        # Word documents
        elif mime in (
            'application/msword',
            'application/vnd.openxmlformats-officedocument.wordprocessingml.document',
        ) or filepath.lower().endswith(('.doc', '.docx')):
            try:
                from docx import Document
                doc = Document(filepath)
                texts = []
                for p in doc.paragraphs[:50]:
                    if p.text:
                        texts.append(p.text)
                snippet = '\n'.join(texts)[:1200]
            except Exception:
                snippet = None

        # PowerPoint
        elif mime in (
            'application/vnd.ms-powerpoint',
            'application/vnd.openxmlformats-officedocument.presentationml.presentation',
        ) or filepath.lower().endswith(('.ppt', '.pptx')):
            try:
                from pptx import Presentation
                prs = Presentation(filepath)
                texts = []
                for i, slide in enumerate(prs.slides[:10]):
                    for shape in slide.shapes:
                        try:
                            if hasattr(shape, 'text') and shape.text:
                                texts.append(shape.text)
                        except Exception:
                            pass
                snippet = '\n'.join(texts)[:1200]
            except Exception:
                snippet = None

        # PDFs
        elif mime == 'application/pdf' or filepath.lower().endswith('.pdf'):
            try:
                import PyPDF2
                with open(filepath, 'rb') as f:
                    reader = PyPDF2.PdfReader(f)
                    texts = []
                    for i, page in enumerate(reader.pages[:5]):
                        try:
                            texts.append(page.extract_text() or '')
                        except Exception:
                            pass
                    snippet = '\n'.join(texts)[:1200]
            except Exception:
                snippet = None

        # Images with optional OCR
        elif mime.startswith('image/') or filepath.lower().endswith(('.png', '.jpg', '.jpeg', '.gif', '.svg')):
            try:
                import pytesseract
                from PIL import Image
                if filepath.lower().endswith('.svg'):
                    with open(filepath, 'r', encoding='utf-8', errors='ignore') as f:
                        snippet = f.read(1200)
                else:
                    img = Image.open(filepath)
                    text = pytesseract.image_to_string(img)
                    snippet = text[:1200]
            except Exception:
                snippet = None

        # Archives: list contained file names; avoid extraction for safety
        elif mime in ('application/zip', 'application/x-zip-compressed') or filepath.lower().endswith('.zip'):
            try:
                import zipfile
                with zipfile.ZipFile(filepath) as z:
                    names = z.namelist()[:20]
                    snippet = ('\n'.join(names))[:1200]
            except Exception:
                snippet = None
        elif mime in ('application/x-rar-compressed', 'application/vnd.rar') or filepath.lower().endswith('.rar'):
            try:
                import rarfile
                rf = rarfile.RarFile(filepath)
                names = rf.namelist()[:20]
                snippet = ('\n'.join(names))[:1200]
            except Exception:
                snippet = None

        else:
            snippet = None
    except Exception:
        snippet = None
    return snippet

@app.route('/api/upload-any', methods=['POST'])
@roles_required('Admin', 'Faculty')
def upload_any_files():
    # Accept files[] list or any arbitrary multipart fields
    files = request.files.getlist('files')
    if not files:
        # Fallback: collect any FileStorage objects
        files = [f for _, f in request.files.items()]
    if not files:
        return jsonify({'error': 'No files provided'}), 400

    upload_dir = os.path.join(app.root_path, 'static', 'uploads')
    os.makedirs(upload_dir, exist_ok=True)

    results = []
    for f in files:
        try:
            safe_name = secure_filename(f.filename or f.name or f.mimetype or f"file_{datetime.utcnow().timestamp()}")
            if not safe_name:
                safe_name = f"file_{int(datetime.utcnow().timestamp())}"
            target_path = os.path.join(upload_dir, safe_name)
            f.save(target_path)

            mime, size = _detect_mime_and_size(target_path, safe_name, getattr(f, 'mimetype', None))
            snippet = _extract_text_snippet_from_path(target_path, mime)

            status = 'processed' if snippet is not None else 'uploaded'
            error_message = None
            if (mime == 'application/pdf' or safe_name.lower().endswith('.pdf')) and snippet is None:
                status = 'error'
                error_message = 'The PDF file could not be processed. Please upload a clearer/standard text-based PDF.'

            remote_url = _optional_cloud_upload(target_path, safe_name, mime)
            log = UniversalUploadLog(
                filename=safe_name,
                mime_type=mime,
                size_bytes=size or 0,
                storage_path=remote_url or target_path,
                status=status,
                error_message=error_message,
                extracted_text_snippet=snippet,
            )
            db.session.add(log)
            db.session.commit()

            payload = {
                'id': log.id,
                'filename': safe_name,
                'mime_type': mime,
                'size': size,
                'upload_time': log.created_at.isoformat() + 'Z',
                'upload_status': log.status,
                'preview': {
                    'text_excerpt': snippet,
                    'url': (remote_url if (remote_url and mime.startswith('image/')) else (url_for('static', filename=f'uploads/{safe_name}') if mime.startswith('image/') else None))
                }
            }
            if error_message:
                payload['error'] = error_message
            results.append(payload)
        except Exception as e:
            # On failure, log the error and continue
            try:
                log = UniversalUploadLog(
                    filename=(f.filename or 'unknown'),
                    mime_type=getattr(f, 'mimetype', None),
                    size_bytes=None,
                    storage_path=None,
                    status='error',
                    error_message=str(e),
                )
                db.session.add(log)
                db.session.commit()
                results.append({
                    'id': log.id,
                    'filename': log.filename,
                    'mime_type': log.mime_type,
                    'size': None,
                    'upload_time': log.created_at.isoformat() + 'Z',
                    'upload_status': 'error',
                    'error': 'Corrupted or unreadable file',
                })
            except Exception:
                results.append({
                    'filename': (getattr(f, 'filename', None) or 'unknown'),
                    'upload_status': 'error',
                    'error': 'Corrupted or unreadable file',
                })

    return jsonify({'files': results}), 201

# List recent universal upload logs
@app.route('/api/uploads/universal', methods=['GET'])
@roles_required('Admin', 'Faculty')
def list_universal_uploads():
    limit = int(request.args.get('limit', '50'))
    items = UniversalUploadLog.query.order_by(UniversalUploadLog.created_at.desc()).limit(limit).all()
    return jsonify({'items': [i.to_dict() for i in items]})


# API Endpoint: Analytics grouped by faculty/upload
@app.route('/api/analysis/by-faculty', methods=['GET'])
def get_analysis_by_faculty():
    uploads = FacultyUpload.query.order_by(FacultyUpload.analysis_date.desc()).all()

    def categorize(name: str):
        course = (name or '').upper()
        if any(kw in course for kw in ['MOBILE', 'COMPUTING', 'LAB', 'PRACTICUM', 'INTEGRATION']):
            return 'Lab'
        if any(kw in course for kw in ['THEORY', 'PRINCIPLES', 'CONCEPTS', 'WORLD']):
            return 'Theory'
        if any(kw in course for kw in ['SOFTWARE', 'ENGINEERING', 'ALGORITHM', 'COMPLEXITY']):
            return 'Software'
        if any(kw in course for kw in ['APPLICATION', 'DEVELOPMENT', 'MULTIMEDIA']):
            return 'Development'
        if any(kw in course for kw in ['DATABASE', 'WEB', 'CLOUD']):
            return 'Web/DB'
        if any(kw in course for kw in ['ALGORITHM', 'DATA STRUCTURE']):
            return 'Algorithm'
        if any(kw in course for kw in ['ETHICS', 'ASSURANCE', 'SECURITY']):
            return 'Security'
        if any(kw in course for kw in ['COMMUNICATION', 'ART', 'APPRECIATION']):
            return 'Humanities'
        return 'General'

    faculties = []
    seen_idem = set()
    print(f"DEBUG: Found {len(uploads)} total uploads in DB")
    for upload in uploads:
        # Filter out subject accomplishment uploads from this view as they don't contain class analysis stats
        try:
            paths_meta = json.loads(upload.file_paths) if upload.file_paths else {}
            if isinstance(paths_meta, dict) and paths_meta.get('metadata', {}).get('module') == 'subject_accomplishment':
                continue
        except:
            pass

        subjects = SubjectAnalysis.query.filter_by(upload_id=upload.id).all()
        analytics_items = AnalyticsData.query.filter_by(upload_id=upload.id).order_by(AnalyticsData.analysis_date.desc()).all()
        
        # DEBUG LOG
        # print(f"DEBUG: Processing Upload ID {upload.id}, Faculty: {upload.faculty_name}, Subjects: {len(subjects)}")

        # Derive program and semester from stored metadata or subjects
        program_val = ''
        semester_val = ''
        try:
            meta = json.loads(upload.file_paths or '{}')
            program_val = (meta.get('metadata', {}) or {}).get('program') or ''
            semester_val = (meta.get('metadata', {}) or {}).get('semester') or ''
            idem_val = (meta.get('metadata', {}) or {}).get('idempotency_key') or ''
        except Exception:
            pass

        # Skip duplicates with same idempotency key scoped to faculty (keep newest)
        if idem_val:
            scoped = f"{upload.faculty_name}:{idem_val}"
            if scoped in seen_idem:
                continue
            seen_idem.add(scoped)
        if not program_val:
            # Fallback: most common program in subjects
            counts = {}
            for s in subjects:
                p = (s.program or '').strip()
                if not p:
                    continue
                counts[p] = counts.get(p, 0) + 1
            if counts:
                program_val = sorted(counts.items(), key=lambda x: (-x[1], x[0]))[0][0]

        # Deduplicate subjects by course within this upload and ensure canonical totals
        agg_subjects = {}
        for s in subjects:
            key = (s.course or '').strip().upper()
            if not key:
                continue
            rec = agg_subjects.get(key) or {
                'name': (s.course or '').strip(),
                'enrolled': 0,
                'passed': 0,
                'failed': 0,
                'deficiencies': 0,
            }
            rec['enrolled'] += int(s.enrolled or 0)
            rec['passed'] += int(s.passed or 0)
            rec['failed'] += int(s.failed or 0)
            rec['deficiencies'] += int(getattr(s, 'num_def', 0) or 0)
            agg_subjects[key] = rec
        subject_chart = []
        for _, rec in agg_subjects.items():
            total = int(rec['enrolled'] or 0)
            passed = int(rec['passed'] or 0)
            # Clamp and recompute failed
            if passed > total:
                passed = total
            failed = max(total - passed, 0)
            pass_rate = (passed / total * 100.0) if total else 0.0
            subject_chart.append({
                'name': rec['name'],
                'passRate': round(pass_rate, 2),
                'enrolled': total,
                'deficiencies': int(rec['deficiencies'] or 0),
                'failed': failed,
                'passed': passed,
            })

        category_perf = {}
        for s in subjects:
            cat = categorize(s.course)
            if cat not in category_perf:
                category_perf[cat] = {'total_pass_rate': 0.0, 'count': 0}
            category_perf[cat]['total_pass_rate'] += float(s.pass_rate or 0)
            category_perf[cat]['count'] += 1

        category_chart = []
        for cat, data in category_perf.items():
            avg = (data['total_pass_rate'] / data['count']) if data['count'] else 0
            category_chart.append({'name': cat, 'avg_pass_rate': avg})

        total_enrolled = sum(int(s.enrolled or 0) for s in subjects)
        total_passed = sum(int(s.passed or 0) for s in subjects)
        total_failed = sum(int(s.failed or 0) for s in subjects)
        total_def = sum(int(getattr(s, 'num_def', 0) or 0) for s in subjects)
        pass_rate = (total_passed / total_enrolled * 100.0) if total_enrolled else 0.0
        fail_rate = (total_failed / total_enrolled * 100.0) if total_enrolled else 0.0

        processed_analytics = []
        detailed_output = ''
        for item in analytics_items:
            try:
                content = json.loads(item.data_content)
            except json.JSONDecodeError:
                content = item.data_content
            processed_analytics.append({
                'id': item.id,
                'upload_id': item.upload_id,
                'data_type': item.data_type,
                'content': content,
                'analysis_date': item.analysis_date.isoformat(),
            })
            if item.data_type == 'detailed_analytics_output' and not detailed_output:
                detailed_output = item.data_content

        # Determine academic year: Explicit column > Metadata > Analysis Date
        display_year = upload.academic_year
        if not display_year:
            try:
                m_meta = json.loads(upload.file_paths or '{}').get('metadata', {})
                display_year = m_meta.get('academic_year') or m_meta.get('year')
            except:
                pass
        if not display_year:
             display_year = str((upload.analysis_date.replace(tzinfo=UTC_TZ)).astimezone(MANILA_TZ).year)

        faculties.append({
            'upload_id': upload.id,
            'faculty_name': upload.faculty_name,
            'analysis_date': to_manila_iso(upload.analysis_date),
            'program': program_val,
            'semester': semester_val,
            'year': str(display_year),
            'subjects': subject_chart,
            'categories': category_chart,
            'summary': {
                'enrolled': total_enrolled,
                'passed': total_passed,
                'failed': total_failed,
                'deficiencies': total_def,
                'pass_rate': round(pass_rate, 2),
                'fail_rate': round(fail_rate, 2),
            },
            'analytics': processed_analytics,
            'detailed_output': detailed_output,
        })

    return jsonify({'faculties': faculties})

# API Endpoint: Get Report Filter Options
@app.route('/api/reports/options', methods=['GET'])
def get_report_options():
    try:
        # Get unique subjects and faculties from uploads/analysis
        subjects = set()
        faculties = set()
        years = set()
        
        uploads = FacultyUpload.query.all()
        for u in uploads:
            faculties.add(u.faculty_name)
            
            # Determine academic year
            display_year = u.academic_year
            if not display_year:
                try:
                    m_meta = json.loads(u.file_paths or '{}').get('metadata', {})
                    display_year = m_meta.get('academic_year') or m_meta.get('year')
                except:
                    pass
            if not display_year and u.analysis_date:
                 display_year = str((u.analysis_date.replace(tzinfo=UTC_TZ)).astimezone(MANILA_TZ).year)
            
            if display_year:
                years.add(str(display_year))
            
            # For subjects, we need to query SubjectAnalysis
            upload_subjects = SubjectAnalysis.query.filter_by(upload_id=u.id).all()
            for sub in upload_subjects:
                subjects.add(sub.course)
                
        return jsonify({
            'subjects': sorted(list(subjects)),
            'faculties': sorted(list(faculties)),
            'years': sorted(list(years), reverse=True)
        })
    except Exception as e:
        return jsonify({'error': str(e)}), 500

# API Endpoint: Generate Report Preview Data
@app.route('/api/reports/preview', methods=['GET'])
def preview_report():
    year_filter = request.args.get('year')
    sem_filter = request.args.get('semester')
    prog_filter = request.args.get('program')
    subj_filter = request.args.get('subject')
    faculty_filter = request.args.get('faculty')
    
    data = _generate_report_data(year_filter, sem_filter, prog_filter, subj_filter, faculty_filter)
    return jsonify(data)

def _generate_report_data(year_filter, sem_filter, prog_filter, subj_filter, faculty_filter):
    uploads = FacultyUpload.query.all()
    filtered_upload_ids = []
    stats = {
        'total_subjects': set(),
        'total_faculty': set(),
        'enrolled': 0,
        'passed': 0,
        'failed': 0,
        'deficiencies': 0,
        'bscs_stats': {'enrolled': 0, 'passed': 0, 'failed': 0, 'deficiencies': 0, 'count': 0, 'subjects': []},
        'bsit_stats': {'enrolled': 0, 'passed': 0, 'failed': 0, 'deficiencies': 0, 'count': 0, 'subjects': []}
    }
    faculty_stats_map = {} 
    subject_details = []
    internal_reviews = []
    mappings = _load_career_dataset() or {}
    skills_map = mappings.get('skills_by_course', {})
    careers_map = mappings.get('careers_by_skill', {})
    accumulated_skills = set()
    risk_distribution = {'Excellent': 0, 'Satisfactory': 0, 'At Risk': 0}
    pass_rate_distribution = {'<75%': 0, '75-80%': 0, '80-85%': 0, '85-90%': 0, '90-100%': 0}
    detailed_failed_students = []
    detailed_risk_students = []
    for u in uploads:
        try:
            meta = json.loads(u.file_paths or '{}')
            m = meta.get('metadata', {})
            u_sem = m.get('semester')
            
            # Prioritize explicit column
            u_year = u.academic_year
            if not u_year:
                u_year = m.get('academic_year') or m.get('year')
            
            if not u_year and u.analysis_date:
                 u_year = str((u.analysis_date.replace(tzinfo=UTC_TZ)).astimezone(MANILA_TZ).year)
            if year_filter and year_filter != 'All':
                if year_filter not in str(u_year or ''):
                    continue
            if sem_filter and sem_filter != 'All':
                if sem_filter.lower() not in str(u_sem or '').lower():
                    continue
            if prog_filter and prog_filter != 'All':
                u_prog = m.get('program')
                if u_prog and prog_filter.lower() not in u_prog.lower():
                    continue
            if faculty_filter and faculty_filter != 'All':
                if faculty_filter.lower() not in u.faculty_name.lower():
                    continue
            filtered_upload_ids.append(u.id)
            stats['total_faculty'].add(u.faculty_name)
            if u.faculty_name not in faculty_stats_map:
                faculty_stats_map[u.faculty_name] = {'enrolled': 0, 'passed': 0, 'deficiencies': 0, 'subjects': 0}
            
            # Parse Deficiency Report if available (DOCX)
            files_meta = meta.get('files', {})
            def_report_path = files_meta.get('def_report')
            if def_report_path and os.path.exists(def_report_path):
                try:
                    doc = Document(def_report_path)
                    for table in doc.tables:
                        for row in table.rows:
                            cells = [c.text.strip() for c in row.cells]
                            if not cells: continue
                            # Check if header or footer
                            if 'COURSE' in cells[0].upper() or 'PREPARED BY' in cells[0].upper():
                                continue
                            if len(cells) >= 4:
                                # Expected: COURSE | PROG | NAME | GRADE | REASON
                                # Map columns based on observed structure
                                subj_code = cells[0]
                                prog_section = cells[1] if len(cells) > 1 else 'N/A'
                                stud_name = cells[2] if len(cells) > 2 else 'Unknown'
                                grade = cells[3] if len(cells) > 3 else 'N/A'
                                
                                if stud_name and stud_name.lower() != 'name of student':
                                    detailed_failed_students.append({
                                        'name': stud_name,
                                        'subject': subj_code,
                                        'program': prog_section,
                                        'grade': grade,
                                        'faculty': u.faculty_name
                                    })
                except Exception as e:
                    print(f"Error reading deficiency report for {u.faculty_name}: {e}")

            subjects = SubjectAnalysis.query.filter_by(upload_id=u.id).all()
            for sub in subjects:
                sub_prog = (sub.program or '').upper()
                if prog_filter and prog_filter != 'All' and prog_filter not in sub_prog:
                    continue
                if subj_filter and subj_filter != 'All' and subj_filter.lower() not in sub.course.lower():
                    continue
                stats['total_subjects'].add(sub.course)
                stats['enrolled'] += sub.enrolled
                stats['passed'] += sub.passed
                stats['failed'] += sub.failed
                stats['deficiencies'] += (sub.num_def or 0)
                faculty_stats_map[u.faculty_name]['enrolled'] += sub.enrolled
                faculty_stats_map[u.faculty_name]['passed'] += sub.passed
                faculty_stats_map[u.faculty_name]['deficiencies'] += (sub.num_def or 0)
                faculty_stats_map[u.faculty_name]['subjects'] += 1
                if 'BSCS' in sub_prog:
                    stats['bscs_stats']['enrolled'] += sub.enrolled
                    stats['bscs_stats']['passed'] += sub.passed
                    stats['bscs_stats']['failed'] += sub.failed
                    stats['bscs_stats']['deficiencies'] += (sub.num_def or 0)
                    stats['bscs_stats']['count'] += 1
                    stats['bscs_stats']['subjects'].append(sub)
                elif 'BSIT' in sub_prog:
                    stats['bsit_stats']['enrolled'] += sub.enrolled
                    stats['bsit_stats']['passed'] += sub.passed
                    stats['bsit_stats']['failed'] += sub.failed
                    stats['bsit_stats']['deficiencies'] += (sub.num_def or 0)
                    stats['bsit_stats']['count'] += 1
                    stats['bsit_stats']['subjects'].append(sub)
                pass_rate = sub.pass_rate
                
                # Update Distributions
                if pass_rate >= 90: pass_rate_distribution['90-100%'] += 1
                elif pass_rate >= 85: pass_rate_distribution['85-90%'] += 1
                elif pass_rate >= 80: pass_rate_distribution['80-85%'] += 1
                elif pass_rate >= 75: pass_rate_distribution['75-80%'] += 1
                else: pass_rate_distribution['<75%'] += 1
                
                # Update Risk Distribution
                if pass_rate >= 90: risk_distribution['Excellent'] += 1
                elif pass_rate >= 75: risk_distribution['Satisfactory'] += 1
                else: risk_distribution['At Risk'] += 1

                badge = 'Satisfactory'
                if pass_rate >= 95: badge = 'Excellent'
                elif pass_rate < 90: badge = 'Needs Improvement'
                subject_details.append({
                    'code': sub.course,
                    'title': sub.course, 
                    'program': sub.program,
                    'enrolled': sub.enrolled,
                    'passed': sub.passed,
                    'failed': sub.failed,
                    'pass_rate': round(sub.pass_rate, 2),
                    'deficiencies': sub.num_def,
                    'badge': badge,
                    'faculty': u.faculty_name
                })
                accomplishments = SubjectAccomplishment.query.filter_by(subject_analysis_id=sub.id).all()
                if not accomplishments:
                     accomplishments = SubjectAccomplishment.query.filter_by(upload_id=u.id, subject_code=sub.course).all()
                for acc in accomplishments:
                    internal_reviews.append({
                        'subject': sub.course,
                        'weakness': acc.weakness,
                        'action': acc.action_taken,
                        'recommendation': acc.recommendation,
                        'faculty': u.faculty_name
                    })
        except Exception as e:
            print(f"Error processing upload {u.id}: {e}")
            continue
    total_unique_students = 0
    if filtered_upload_ids:
        profile_count = StudentAcademicProfile.query.count()
        if profile_count > 0:
             total_unique_students = profile_count
        else:
             total_unique_students = stats['enrolled']
    overall_pass_rate = (stats['passed'] / stats['enrolled'] * 100) if stats['enrolled'] else 0
    overall_fail_rate = (stats['failed'] / stats['enrolled'] * 100) if stats['enrolled'] else 0
    consistency_flag = False
    consistency_msg = ""
    if abs(stats['enrolled'] - (stats['passed'] + stats['failed'])) > (stats['enrolled'] * 0.1): 
         consistency_flag = True
         consistency_msg = f"Data mismatch: Enrolled ({stats['enrolled']}) vs Outcomes ({stats['passed'] + stats['failed']})"
    interpretation = 'Satisfactory'
    if overall_pass_rate >= 95: interpretation = 'Excellent'
    elif overall_pass_rate < 90: interpretation = 'Needs Improvement'
    def get_program_extremes(subjects_list):
        if not subjects_list:
            return "N/A", "N/A"
        sorted_subs = sorted(subjects_list, key=lambda s: s.pass_rate)
        weakest = f"{sorted_subs[0].course} ({round(sorted_subs[0].pass_rate, 2)}%)"
        strongest = f"{sorted_subs[-1].course} ({round(sorted_subs[-1].pass_rate, 2)}%)"
        return strongest, weakest
    bscs_stats = stats['bscs_stats']
    bsit_stats = stats['bsit_stats']
    bscs_pr = (bscs_stats['passed'] / bscs_stats['enrolled'] * 100) if bscs_stats['enrolled'] else 0
    bsit_pr = (bsit_stats['passed'] / bsit_stats['enrolled'] * 100) if bsit_stats['enrolled'] else 0
    bscs_strongest, bscs_weakest = get_program_extremes(bscs_stats['subjects'])
    bsit_strongest, bsit_weakest = get_program_extremes(bsit_stats['subjects'])
    faculty_performance = []
    for fname, fstats in faculty_stats_map.items():
        if fstats['enrolled'] == 0: continue
        fpr = (fstats['passed'] / fstats['enrolled'] * 100)
        faculty_performance.append({
            'name': fname,
            'students': fstats['enrolled'],
            'pass_rate': round(fpr, 2),
            'deficiencies': fstats['deficiencies']
        })
    hotspots = []
    top_def_subs = sorted(subject_details, key=lambda x: x['deficiencies'], reverse=True)
    for sub in top_def_subs:
        if sub['deficiencies'] >= 5: 
            exists = next((h for h in hotspots if h['subject'] == sub['code']), None)
            if not exists:
                hotspots.append({
                    'subject': sub['code'],
                    'reason': f"{sub['deficiencies']} Incomplete counts (High volume)"
                })
            if len(hotspots) >= 3: break
    low_pass_subs = sorted(subject_details, key=lambda x: x['pass_rate'])
    for sub in low_pass_subs:
        if sub['pass_rate'] < 85: 
            exists = next((h for h in hotspots if h['subject'] == sub['code']), None)
            if not exists:
                hotspots.append({
                    'subject': sub['code'],
                    'reason': f"Critical pass rate of {sub['pass_rate']}%"
                })
                if len(hotspots) >= 5: break
    if not hotspots:
        hotspots.append({'subject': 'None', 'reason': 'No critical deficiency hotspots identified for this period.'})
    intervention_list = []
    high_risk_profiles = StudentAcademicProfile.query.filter(StudentAcademicProfile.risk_level.in_(['High', 'Medium'])).all()
    if high_risk_profiles:
        for p in high_risk_profiles:
            masked_id = 'N/A'
            if p.student_id:
                if len(p.student_id) > 4:
                    masked_id = '*' * (len(p.student_id) - 4) + p.student_id[-4:]
                else:
                    masked_id = p.student_id
            program_val = 'N/A'
            if p.student_id:
                u_rec = User.query.filter_by(student_id=p.student_id).first()
                if u_rec and u_rec.program:
                    program_val = u_rec.program
            intervention_list.append({
                'name': p.student_name,
                'id': masked_id,
                'program': program_val,
                'risk': p.risk_level,
                'details': f"Failed: {p.failed_count}, Avg: {p.average_grade}"
            })
    else:
        from collections import defaultdict
        student_failures = defaultdict(list)
        student_grades = defaultdict(list)
        student_programs = {}
        
        for fail in detailed_failed_students:
            student_failures[fail['name']].append(fail['subject'])
            student_grades[fail['name']].append(fail['grade'])
            if fail.get('program') and fail['program'] != 'N/A':
                student_programs[fail['name']] = fail['program']
        
        candidates = []
        for name, subs in student_failures.items():
            grades = student_grades[name]
            count = len(subs)
            has_5 = any('5.0' in str(g) or '5.00' in str(g) for g in grades)
            
            risk = 'Medium'
            if count >= 2 or has_5:
                risk = 'High'
            
            score = count * 10 + (5 if has_5 else 0)
            
            candidates.append({
                'name': name,
                'count': count,
                'subs': subs,
                'risk': risk,
                'score': score,
                'program': student_programs.get(name, 'N/A')
            })
        
        # Sort by score desc (Highest risk first)
        candidates.sort(key=lambda x: x['score'], reverse=True)
        
        for cand in candidates[:10]:
            masked_id = 'N/A'
            program_val = cand['program']
            
            # Try to lookup User for ID if missing
            u_rec = User.query.filter_by(full_name=cand['name']).first()
            if u_rec:
                if u_rec.student_id:
                    sid = u_rec.student_id
                    if len(sid) > 4:
                        masked_id = '*' * (len(sid) - 4) + sid[-4:]
                    else:
                        masked_id = sid
                if program_val == 'N/A' and u_rec.program:
                    program_val = u_rec.program
            
            intervention_list.append({
                'name': cand['name'],
                'id': masked_id,
                'program': program_val,
                'risk': cand['risk'],
                'details': f"Failed {cand['count']} subjects: {', '.join(cand['subs'])}"
            })
    intervention_list = intervention_list[:10]
    narrative_parts = []
    if hotspots and hotspots[0]['subject'] != 'None':
        narrative_parts.append(f"Immediate academic intervention is recommended for {hotspots[0]['subject']} due to {hotspots[0]['reason'].lower()}.")
    if bscs_pr < 80 or bsit_pr < 80:
         narrative_parts.append("A review of curriculum delivery is suggested for programs with pass rates below 80%.")
    if intervention_list:
        narrative_parts.append(f"The system has identified {len(intervention_list)} students requiring prioritized academic advising.")
    if not narrative_parts:
        narrative_parts.append("Academic performance is stable. Continue monitoring student progress and faculty resource allocation.")
    narrative = " ".join(narrative_parts)
    response_data = {
        'header': {
            'year': year_filter or '2024-2025',
            'semester': sem_filter or 'All',
            'program': prog_filter,
            'subject': subj_filter,
            'faculty': faculty_filter,
            'generated_at': to_manila_iso(datetime.now())
        },
        'executive_summary': {
            'total_subjects': len(stats['total_subjects']),
            'total_faculty': len(stats['total_faculty']),
            'total_students': total_unique_students,
            'overall_pass_rate': round(overall_pass_rate, 2),
            'overall_fail_rate': round(overall_fail_rate, 2),
            'total_deficiencies': stats['deficiencies'],
            'interpretation': interpretation,
            'data_integrity_warning': consistency_msg if consistency_flag else None
        },
        'distributions': {
            'risk': risk_distribution,
            'pass_rate': pass_rate_distribution
        },
        'program_comparison': {
            'bscs': {
                'passed': bscs_stats['passed'],
                'failed': bscs_stats['failed'],
                'pass_rate': round(bscs_pr, 2),
                'deficiencies': bscs_stats['deficiencies'],
                'strongest': bscs_strongest,
                'weakest': bscs_weakest
            },
            'bsit': {
                'passed': bsit_stats['passed'],
                'failed': bsit_stats['failed'],
                'pass_rate': round(bsit_pr, 2),
                'deficiencies': bsit_stats['deficiencies'],
                'strongest': bsit_strongest,
                'weakest': bsit_weakest
            }
        },
        'faculty_performance': faculty_performance,
        'subject_performance': subject_details,
        'internal_reviews': internal_reviews,
        'hotspots': hotspots,
        'student_analysis': {
            'intervention_list': intervention_list
        },
        'narrative': narrative
    }
    return response_data

@app.route('/api/reports/export/docx', methods=['GET'])
def export_report_docx():
    year_filter = request.args.get('year')
    sem_filter = request.args.get('semester')
    prog_filter = request.args.get('program')
    subj_filter = request.args.get('subject')
    faculty_filter = request.args.get('faculty')
    data = _generate_report_data(year_filter, sem_filter, prog_filter, subj_filter, faculty_filter)
    document = Document()
    style = document.styles['Normal']
    font = style.font
    font.name = 'Times New Roman'
    font.size = Pt(12)

    # --- Header with Logos ---
    header_table = document.add_table(rows=1, cols=3)
    header_table.autofit = False
    # Set widths: Logo(1.0) | Text(4.5) | Logo(1.0)
    # Note: Column width setting in python-docx can be tricky, relying on cell width often works better or just letting it autofit if content dictates.
    # We'll try to set cell widths explicitly if needed, but adding pictures usually constrains it.
    
    # Left Logo (LSPU)
    cell_left = header_table.cell(0, 0)
    p_left = cell_left.paragraphs[0]
    p_left.alignment = WD_ALIGN_PARAGRAPH.CENTER
    try:
        lspu_logo_path = os.path.join(project_root, 'frontend', 'public', 'lspu.jfif')
        if os.path.exists(lspu_logo_path):
            run_left = p_left.add_run()
            run_left.add_picture(lspu_logo_path, width=Inches(0.9))
    except Exception as e:
        print(f"Error loading LSPU logo: {e}")

    # Center Text
    cell_center = header_table.cell(0, 1)
    p_center = cell_center.paragraphs[0]
    p_center.alignment = WD_ALIGN_PARAGRAPH.CENTER
    
    r1 = p_center.add_run('Republic of the Philippines\n')
    r1.font.size = Pt(10)
    
    r2 = p_center.add_run('Laguna State Polytechnic University\n')
    r2.font.name = 'Old English Text MT'
    r2.font.size = Pt(16)
    r2.bold = True
    
    r3 = p_center.add_run('Province of Laguna\n')
    r3.font.size = Pt(10)
    
    r4 = p_center.add_run('COLLEGE OF COMPUTER STUDIES')
    r4.font.name = 'Arial'
    r4.font.size = Pt(14)
    r4.bold = True

    # Right Logo (CCS)
    cell_right = header_table.cell(0, 2)
    p_right = cell_right.paragraphs[0]
    p_right.alignment = WD_ALIGN_PARAGRAPH.CENTER
    try:
        ccs_logo_path = os.path.join(project_root, 'frontend', 'public', 'lspui-seal-ccs.jfif')
        if os.path.exists(ccs_logo_path):
            run_right = p_right.add_run()
            run_right.add_picture(ccs_logo_path, width=Inches(0.9))
        else:
            print(f"CCS logo not found at: {ccs_logo_path}")
    except Exception as e:
        print(f"Error loading CCS logo: {e}")

    # Meta Info
    p3 = document.add_paragraph(f"Academic Year: {data['header']['year']} | Semester: {data['header']['semester']}")
    p3.alignment = WD_ALIGN_PARAGRAPH.CENTER
    document.add_heading('EXECUTIVE SUMMARY', level=1)
    es = data['executive_summary']
    if es['data_integrity_warning']:
        p_warn = document.add_paragraph()
        run_warn = p_warn.add_run(f"WARNING: {es['data_integrity_warning']}")
        run_warn.font.color.rgb = RGBColor(255, 0, 0)
        run_warn.bold = True
    table = document.add_table(rows=1, cols=2)
    table.style = 'Table Grid'
    hdr_cells = table.rows[0].cells
    hdr_cells[0].text = 'Metric'
    hdr_cells[1].text = 'Value'
    metrics = [
        ('Total Students', es['total_students']),
        ('Overall Pass Rate', f"{es['overall_pass_rate']}%"),
        ('Overall Fail Rate', f"{es['overall_fail_rate']}%"),
        ('Analyzed Courses', es['total_subjects']),
        ('Faculty Count', es['total_faculty']),
        ('Total Deficiencies', es['total_deficiencies']),
        ('Overall Interpretation', es['interpretation'])
    ]
    for m, v in metrics:
        row_cells = table.add_row().cells
        row_cells[0].text = str(m)
        row_cells[1].text = str(v)
    document.add_heading('PROGRAM PERFORMANCE ANALYSIS', level=1)
    pc = data['program_comparison']
    table_prog = document.add_table(rows=1, cols=4)
    table_prog.style = 'Table Grid'
    hdr_prog = table_prog.rows[0].cells
    hdr_prog[0].text = 'Program'
    hdr_prog[1].text = 'Pass Rate'
    hdr_prog[2].text = 'Strongest Subject'
    hdr_prog[3].text = 'Weakest Subject'
    for prog_name, pdata in pc.items():
        if pdata['pass_rate'] == 0 and pdata['passed'] == 0: continue
        row = table_prog.add_row().cells
        row[0].text = prog_name.upper()
        row[1].text = f"{pdata['pass_rate']}%"
        row[2].text = str(pdata['strongest'])
        row[3].text = str(pdata['weakest'])
    document.add_heading('FACULTY PERFORMANCE OVERVIEW', level=1)
    fp = data['faculty_performance']
    if fp:
        table_fac = document.add_table(rows=1, cols=4)
        table_fac.style = 'Table Grid'
        hf = table_fac.rows[0].cells
        hf[0].text = 'Faculty Name'
        hf[1].text = 'Students'
        hf[2].text = 'Pass Rate'
        hf[3].text = 'Deficiencies'
        for f in fp:
            row = table_fac.add_row().cells
            row[0].text = f['name']
            row[1].text = str(f['students'])
            row[2].text = f"{f['pass_rate']}%"
            row[3].text = str(f['deficiencies'])
    else:
        document.add_paragraph("No faculty data available.")
    document.add_heading('CRITICAL DEFICIENCY HOTSPOTS', level=1)
    hs = data['hotspots']
    if hs:
        for h in hs:
            p = document.add_paragraph()
            if h['subject'] != 'None':
                run = p.add_run(f"{h['subject']}: ")
                run.bold = True
                run.font.color.rgb = RGBColor(255, 0, 0)
            p.add_run(h['reason'])
    else:
         document.add_paragraph("No critical hotspots identified.")
    document.add_heading('STUDENT INTERVENTION LIST', level=1)
    si = data['student_analysis']['intervention_list']
    if si:
        for s in si:
            p = document.add_paragraph()
            run = p.add_run(f"{s['name']} (ID: {s['id']}) - {s['risk']} Risk")
            run.bold = True
            p2 = document.add_paragraph(f"Program: {s['program']} | {s['details']}")
            p2.paragraph_format.left_indent = Inches(0.5)
    else:
        document.add_paragraph("No high-risk students identified.")
    document.add_heading('SYSTEM NARRATIVE RECOMMENDATION', level=1)
    document.add_paragraph(data['narrative'])
    section = document.sections[0]
    footer = section.footer
    p_foot = footer.paragraphs[0] if footer.paragraphs else footer.add_paragraph()
    p_foot.text = f"Generated by AI-Powered CQI System | {datetime.now().strftime('%Y-%m-%d %H:%M:%S')}"
    p_foot.alignment = WD_ALIGN_PARAGRAPH.CENTER
    f = BytesIO()
    document.save(f)
    f.seek(0)
    return send_file(
        f,
        as_attachment=True,
        download_name=f"Academic_Evaluation_Report_{datetime.now().strftime('%Y%m%d')}.docx",
        mimetype='application/vnd.openxmlformats-officedocument.wordprocessingml.document'
    )


@app.route('/api/student-risk/extraction-report/latest', methods=['GET'])
def student_risk_extraction_report_latest():
    try:
        upload = FacultyUpload.query.filter(FacultyUpload.file_paths.contains('student_risk')).order_by(FacultyUpload.analysis_date.desc()).first()
        if not upload:
            return jsonify({'report': {}})
        item = AnalyticsData.query.filter_by(upload_id=upload.id, data_type='student_risk_extraction_report').order_by(AnalyticsData.analysis_date.desc()).first()
        payload = json.loads(item.data_content) if item else {}
        return jsonify({'report': payload})
    except Exception as e:
        return jsonify({'error': str(e)}), 500

if __name__ == '__main__':
    print("Starting Flask app...")
    # Disable the Flask reloader to avoid mid-request resets during development
    try:
        app.run(debug=True, host='0.0.0.0', port=5000, use_reloader=False) # Listen on all interfaces, port 5000
    except Exception as e:
        print(f"Failed to start app: {e}")

